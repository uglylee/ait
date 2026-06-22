"""
LangChain 集成模块
提供完整的 LangChain 功能：LLM、Memory、Chains、Agents、RAG、Tools
"""

import os
import json
from typing import List, Dict, Optional, Any
from langchain_core.language_models import BaseLLM
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from langchain_community.document_loaders import (
    TextLoader, PyPDFLoader, CSVLoader
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings


# ============ 1. 自定义 LLM 包装器（对接我们的多模型接口） ============

class MultiProviderLLM(BaseLLM):
    """通用多模型 LLM 包装器，对接项目已有的提供商"""
    
    provider: str = "agnes"
    model: str = None
    base_url: str = None
    api_key: str = None
    temperature: float = 0.7
    max_tokens: int = 2048
    
    def __init__(self, provider: str = None, **kwargs):
        super().__init__(**kwargs)
        if provider is None:
            from dotenv import load_dotenv
            load_dotenv()
            provider = os.getenv("DEFAULT_PROVIDER", "agnes")
        self.provider = provider
        self._load_config()
    
    def _load_config(self):
        from dotenv import load_dotenv
        load_dotenv()
        
        prefix = self.provider.upper()
        self.api_key = os.getenv(f"{prefix}_API_KEY", "")
        self.base_url = os.getenv(f"{prefix}_BASE_URL", "https://api.deepseek.com")
        self.model = os.getenv(f"{prefix}_MODEL", self._default_model())
    
    def _default_model(self) -> str:
        defaults = {
            "deepseek": "deepseek-chat",
            "mimo": "mimo-auto",
            "agnes": "agnes-2.0-flash",
            "qwen": "qwen-turbo",
            "glm": "glm-4-flash",
            "openai": "gpt-3.5-turbo",
        }
        return defaults.get(self.provider, "deepseek-chat")
    
    def _generate(self, messages: List, stop: List[str] = None, **kwargs) -> str:
        import httpx
        lc_messages = []
        for msg in messages:
            if isinstance(msg, HumanMessage):
                lc_messages.append({"role": "user", "content": msg.content})
            elif isinstance(msg, AIMessage):
                lc_messages.append({"role": "assistant", "content": msg.content})
            elif isinstance(msg, SystemMessage):
                lc_messages.append({"role": "system", "content": msg.content})
            elif isinstance(msg, dict):
                lc_messages.append(msg)
        
        with httpx.Client(timeout=60.0) as client:
            response = client.post(
                f"{self.base_url}/v1/chat/completions",
                headers={"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"},
                json={
                    "model": self.model,
                    "messages": lc_messages,
                    "temperature": kwargs.get("temperature", self.temperature),
                    "max_tokens": kwargs.get("max_tokens", self.max_tokens),
                }
            )
            response.raise_for_status()
            data = response.json()
            return data["choices"][0]["message"]["content"]
    
    @property
    def _llm_type(self) -> str:
        return f"multi-provider-{self.provider}"
    
    @property
    def _identifying_params(self) -> dict:
        return {"provider": self.provider, "model": self.model, "base_url": self.base_url}


# ============ 2. Memory（记忆模块） ============

class ConversationMemory:
    """会话记忆管理"""
    
    def __init__(self, max_history: int = 20):
        self.max_history = max_history
        self.messages: List[Dict] = []
    
    def add_user_message(self, content: str):
        self.messages.append({"role": "user", "content": content})
        self._trim()
    
    def add_ai_message(self, content: str):
        self.messages.append({"role": "assistant", "content": content})
        self._trim()
    
    def add_system_message(self, content: str):
        self.messages.insert(0, {"role": "system", "content": content})
    
    def get_messages(self) -> List[Dict]:
        return self.messages.copy()
    
    def get_lc_messages(self) -> List:
        """转为 LangChain 消息格式"""
        lc_messages = []
        for msg in self.messages:
            if msg["role"] == "user":
                lc_messages.append(HumanMessage(content=msg["content"]))
            elif msg["role"] == "assistant":
                lc_messages.append(AIMessage(content=msg["content"]))
            elif msg["role"] == "system":
                lc_messages.append(SystemMessage(content=msg["content"]))
        return lc_messages
    
    def clear(self):
        self.messages.clear()
    
    def _trim(self):
        if len(self.messages) > self.max_history:
            system_msgs = [m for m in self.messages if m["role"] == "system"]
            other_msgs = [m for m in self.messages if m["role"] != "system"]
            self.messages = system_msgs + other_msgs[-self.max_history:]
    
    def summary(self) -> str:
        """生成对话摘要"""
        lines = []
        for msg in self.messages:
            role = "用户" if msg["role"] == "user" else "AI" if msg["role"] == "assistant" else "系统"
            lines.append(f"{role}: {msg['content'][:100]}...")
        return "\n".join(lines)


# ============ 3. Chains（链式调用） ============

class SimpleChain:
    """简单链：prompt -> llm -> output"""
    
    def __init__(self, llm: MultiProviderLLM, template: str):
        self.llm = llm
        self.prompt = ChatPromptTemplate.from_template(template)
        self.chain = self.prompt | self.llm | StrOutputParser()
    
    def run(self, **kwargs) -> str:
        return self.chain.invoke(kwargs)


class ConversationChain:
    """带记忆的对话链"""
    
    def __init__(self, llm: MultiProviderLLM, system_prompt: str = "", memory: ConversationMemory = None):
        self.llm = llm
        self.memory = memory or ConversationMemory()
        self.system_prompt = system_prompt
        
        if system_prompt:
            self.memory.add_system_message(system_prompt)
    
    def run(self, user_input: str) -> str:
        self.memory.add_user_message(user_input)
        
        messages = self.memory.get_lc_messages()
        response = self.llm._generate(messages)
        
        self.memory.add_ai_message(response)
        return response
    
    def get_history(self) -> List[Dict]:
        return self.memory.get_messages()
    
    def clear(self):
        self.memory.clear()
        if self.system_prompt:
            self.memory.add_system_message(self.system_prompt)


class SequentialChain:
    """顺序链：多个步骤串行执行"""
    
    def __init__(self, steps: List[Dict]):
        """
        steps: [{"name": "step1", "llm": llm, "template": "...", "output_key": "result1"}, ...]
        """
        self.steps = steps
    
    def run(self, initial_input: str) -> Dict[str, str]:
        context = {"input": initial_input}
        for step in self.steps:
            llm = step["llm"]
            template = step["template"]
            prompt = ChatPromptTemplate.from_template(template)
            chain = prompt | llm | StrOutputParser()
            result = chain.invoke(context)
            context[step["output_key"]] = result
            context["input"] = result  # 传递给下一步
        return context


# ============ 4. Agents（智能体） ============

class Tool:
    """工具基类"""
    
    def __init__(self, name: str, description: str):
        self.name = name
        self.description = description
    
    def run(self, input: str) -> str:
        raise NotImplementedError


class CalculatorTool(Tool):
    """计算器工具"""
    
    def __init__(self):
        super().__init__("calculator", "Evaluate a math expression. Input: a valid Python math expression, e.g. '2+3*4', '100/7', '2**10'. Returns the numeric result.")
    
    def run(self, input: str) -> str:
        try:
            result = eval(input)
            return str(result)
        except Exception as e:
            return f"计算错误: {e}"


class OpenAppTool(Tool):
    """打开应用工具"""
    
    def __init__(self):
        super().__init__("open_app", "Open a local application by name. Supported: chrome, edge, notepad, vscode, code, explorer, finder, terminal, cmd, calculator, wechat. Input: app name as a string.")
    
    def run(self, input: str) -> str:
        import subprocess
        import platform
        
        app_map = {
            "chrome": {"windows": "start chrome", "darwin": "open -a 'Google Chrome'", "linux": "google-chrome"},
            "edge": {"windows": "start msedge", "darwin": "open -a 'Microsoft Edge'", "linux": "microsoft-edge"},
            "notepad": {"windows": "notepad", "darwin": "open -a TextEdit", "linux": "gedit"},
            "记事本": {"windows": "notepad", "darwin": "open -a TextEdit", "linux": "gedit"},
            "计算器": {"windows": "calc", "darwin": "open -a Calculator", "linux": "gnome-calculator"},
            "vscode": {"windows": "code", "darwin": "open -a 'Visual Studio Code'", "linux": "code"},
            "code": {"windows": "code", "darwin": "open -a 'Visual Studio Code'", "linux": "code"},
            "explorer": {"windows": "explorer .", "darwin": "open .", "linux": "nautilus ."},
            "finder": {"windows": "open .", "darwin": "open .", "linux": "nautilus ."},
            "terminal": {"windows": "cmd", "darwin": "open -a Terminal", "linux": "gnome-terminal"},
            "cmd": {"windows": "cmd", "darwin": "open -a Terminal", "linux": "gnome-terminal"},
            "wechat": {"windows": "start WeChat", "darwin": "open -a WeChat", "linux": ""},
            "微信": {"windows": "start WeChat", "darwin": "open -a WeChat", "linux": ""},
        }
        
        app_name = input.lower().strip()
        system = platform.system().lower()
        os_key = "windows" if system == "windows" else "darwin" if system == "darwin" else "linux"
        
        if app_name in app_map:
            cmd = app_map[app_name].get(os_key, "")
            if cmd:
                try:
                    subprocess.Popen(cmd, shell=True)
                    return f"已打开 {input}"
                except Exception as e:
                    return f"打开失败: {e}"
            return f"不支持的操作系统: {system}"
        
        # 尝试直接运行
        try:
            subprocess.Popen(app_name, shell=True)
            return f"已尝试打开 {input}"
        except Exception as e:
            return f"未识别的应用 '{input}'，打开失败: {e}"


class OpenUrlTool(Tool):
    """打开网页工具"""
    
    def __init__(self):
        super().__init__("open_url", "Open a URL in the default browser or open a local file with its associated application. Input: a URL (http/https) or a local file path.")
    
    def run(self, input: str) -> str:
        import subprocess
        import platform
        import os
        
        target = input.strip()
        
        # 处理 file:/// URL
        if target.startswith("file:///"):
            target = target[8:] if platform.system() == "Windows" else target[7:]
            target = target.replace('/', '\\') if platform.system() == "Windows" else target
        
        # 检查文件是否存在
        if not target.startswith("http") and not os.path.exists(target):
            return f"文件不存在: {target}"
        
        system = platform.system().lower()
        try:
            if system == "windows":
                os.startfile(target)
            elif system == "darwin":
                subprocess.Popen(['open', target])
            else:
                subprocess.Popen(['xdg-open', target])
            return f"已打开: {target}"
        except Exception as e:
            return f"打开失败: {e}"


class OpenFileTool(Tool):
    """打开文件工具"""
    
    def __init__(self):
        super().__init__("open_file", "Open a local file using the OS default application. Input: absolute file path, e.g. 'C:\\Users\\test.txt' or '/home/user/doc.pdf'.")
    
    def run(self, input: str) -> str:
        import os
        import platform
        import subprocess
        
        if not os.path.exists(input):
            return f"文件不存在: {input}"
        
        system = platform.system().lower()
        try:
            if system == "windows":
                os.startfile(input)
            elif system == "darwin":
                subprocess.Popen(['open', input])
            else:
                subprocess.Popen(['xdg-open', input])
            return f"已打开文件: {input}"
        except Exception as e:
            return f"打开失败: {e}"


class RunCommandTool(Tool):
    """执行系统命令工具"""
    
    def __init__(self):
        super().__init__("run_command", "Execute a shell command and return its stdout/stderr. Input: a valid shell command string, e.g. 'ipconfig', 'dir C:\\', 'ls -la', 'python --version', 'pip install requests'. Timeout: 30 seconds.")
    
    def run(self, input: str) -> str:
        import subprocess
        try:
            result = subprocess.run(input, shell=True, capture_output=True, text=True, timeout=30)
            output = result.stdout or result.stderr or "命令执行成功（无输出）"
            return output[:2000]
        except subprocess.TimeoutExpired:
            return "命令执行超时（30秒）"
        except Exception as e:
            return f"执行失败: {e}"


class CodeGeneratorTool(Tool):
    """代码生成 + 运行工具"""
    
    def __init__(self):
        super().__init__("code_generate", "Generate a code project from a description and run it. Input format: 'project_name|description|language'. Language options: python, html, js/javascript. Example: 'my_app|A Flask web server|python'. The tool creates the project files and attempts to run them.")
    
    def run(self, input: str) -> str:
        import os
        import subprocess
        import platform
        
        parts = input.split("|")
        if len(parts) < 2:
            return "格式错误，请输入: 项目名|描述|语言(可选)"
        
        project_name = parts[0].strip()
        description = parts[1].strip()
        lang = parts[2].strip() if len(parts) > 2 else "python"
        
        projects_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_projects")
        project_dir = os.path.join(projects_dir, project_name)
        os.makedirs(project_dir, exist_ok=True)
        
        code = self._generate_code(description, lang)
        
        if lang.lower() in ["python", "py"]:
            file_path = os.path.join(project_dir, "main.py")
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(code)
            
            # 创建说明文件
            with open(os.path.join(project_dir, "README.md"), "w", encoding="utf-8") as f:
                f.write(f"# {project_name}\n\n{description}\n\n## 运行方式\n\n```bash\npython main.py\n```\n")
            
            # 尝试运行
            try:
                system = platform.system().lower()
                cmd = f'cd "{project_dir}" && python main.py'
                if system != "windows":
                    cmd = f'cd "{project_dir}" && python3 main.py'
                
                result = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=30)
                output = result.stdout or ""
                errors = result.stderr or ""
                
                run_result = ""
                if output:
                    run_result = f"\n\n--- 运行结果 ---\n{output[:3000]}"
                if errors and "Error" in errors:
                    run_result += f"\n\n--- 错误信息 ---\n{errors[:1000]}"
                
                return f"项目已生成: {project_dir}\n\n文件: main.py, README.md\n\n--- 代码预览 ---\n```python\n{code[:1500]}\n```{run_result}"
            except subprocess.TimeoutExpired:
                return f"项目已生成: {project_dir}\n\n--- 代码预览 ---\n```python\n{code[:2000]}\n```\n\n运行超时（30秒）"
            except Exception as e:
                return f"项目已生成: {project_dir}\n\n--- 代码预览 ---\n```python\n{code[:2000]}\n```\n\n运行错误: {e}"
        
        elif lang.lower() in ["html", "web"]:
            file_path = os.path.join(project_dir, "index.html")
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(code)
            
            with open(os.path.join(project_dir, "README.md"), "w", encoding="utf-8") as f:
                f.write(f"# {project_name}\n\n{description}\n\n## 运行方式\n\n直接打开 index.html\n")
            
            return f"项目已生成: {file_path}\n\n请用 open_url 工具打开此文件"
        
        elif lang.lower() in ["js", "javascript", "node"]:
            file_path = os.path.join(project_dir, "main.js")
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(code)
            
            with open(os.path.join(project_dir, "package.json"), "w", encoding="utf-8") as f:
                f.write(json.dumps({"name": project_name, "version": "1.0.0", "main": "main.js"}, indent=2))
            
            return f"项目已生成: {project_dir}\n\n文件: main.js, package.json\n\n--- 代码预览 ---\n```javascript\n{code[:1500]}\n```\n\n运行: cd {project_dir} && node main.js"
        
        else:
            file_path = os.path.join(project_dir, f"main.{lang}")
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(code)
            
            return f"项目已生成: {project_dir}\n\n文件: main.{lang}\n\n--- 代码预览 ---\n```\n{code[:1500]}\n```"
    
    def _generate_code(self, description: str, lang: str) -> str:
        """调用LLM生成真实代码"""
        lang_map = {"python": "Python", "py": "Python", "html": "HTML/CSS/JavaScript", "web": "HTML/CSS/JavaScript", "js": "JavaScript", "javascript": "JavaScript", "node": "Node.js"}
        lang_name = lang_map.get(lang.lower(), lang)

        prompt = f"""你是一个高级程序员。请根据以下描述生成一个完整可运行的{lang_name}项目。

描述: {description}

要求:
1. 代码必须完整、可直接运行，不要省略任何部分
2. 代码要美观、功能完整
3. 只输出代码，不要任何解释文字
4. 不要用markdown代码块包裹

代码:"""

        try:
            from llm_config import get_llm_router
            router = get_llm_router()
            llm = router.get_provider(None)
            import httpx
            import json as _json

            with httpx.Client(timeout=120.0) as client:
                resp = client.post(
                    llm.chat_endpoint,
                    headers={"Authorization": f"Bearer {llm.api_key}", "Content-Type": "application/json"},
                    json={
                        "model": llm.default_model,
                        "messages": [{"role": "user", "content": prompt}],
                        "temperature": 0.3,
                        "max_tokens": 8192,
                    }
                )
                data = resp.json()
                code = data["choices"][0]["message"]["content"]

            code = code.strip()
            if code.startswith("```"):
                lines = code.split("\n")
                lines = lines[1:]
                if lines and lines[-1].strip() == "```":
                    lines = lines[:-1]
                code = "\n".join(lines)

            if len(code) < 50:
                return self._fallback_code(description, lang)

            return code
        except Exception as e:
            print(f"[CodeGenerator] LLM call failed: {e}, using fallback")
            return self._fallback_code(description, lang)
    
    def _fallback_code(self, description: str, lang: str) -> str:
        """LLM失败时的模板代码"""
        if lang.lower() in ["python", "py"]:
            return f'''"""
{description}
"""
import sys

def main():
    print("=" * 50)
    print(f"项目: {description}")
    print("=" * 50)
    print()
    print("这是一个示例项目框架。")
    print("请根据需求修改 main.py 文件。")
    print()
    print("功能:")
    print("1. 基础项目结构")
    print("2. 可扩展的模块设计")
    print("3. 错误处理机制")
    print()
    print("运行方式: python main.py")
    print()

if __name__ == "__main__":
    main()
'''
        elif lang.lower() in ["html", "web"]:
            return f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{description}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: -apple-system, sans-serif; background: #f5f5f5; padding: 20px; }}
        .container {{ max-width: 800px; margin: 0 auto; background: #fff; border-radius: 12px; padding: 30px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }}
        h1 {{ color: #333; margin-bottom: 20px; }}
        p {{ color: #666; line-height: 1.8; }}
        .btn {{ display: inline-block; padding: 10px 20px; background: #409eff; color: #fff; border: none; border-radius: 6px; cursor: pointer; margin-top: 15px; }}
        .btn:hover {{ background: #337ecc; }}
    </style>
</head>
<body>
    <div class="container">
        <h1>{description}</h1>
        <p>这是一个自动生成的网页项目。</p>
        <p>你可以修改此文件来实现具体功能。</p>
        <button class="btn" onclick="alert('Hello!')">点击测试</button>
    </div>
</body>
</html>'''
        elif lang.lower() in ["js", "javascript", "node"]:
            return f'''/**
 * {description}
 */

const http = require('http');
const fs = require('fs');
const path = require('path');

const PORT = 3000;

const server = http.createServer((req, res) => {{
    res.writeHead(200, {{ 'Content-Type': 'text/html; charset=utf-8' }});
    res.end(`
        <!DOCTYPE html>
        <html>
        <head><title>{description}</title></head>
        <body>
            <h1>{description}</h1>
            <p>Node.js 服务运行中...</p>
        </body>
        </html>
    `);
}});

server.listen(PORT, () => {{
    console.log(`服务已启动: http://localhost:${{PORT}}`);
}});
'''
        return f'print("{description}")'

    def _general_py(self, description: str) -> str:
        return f'''"""
{description}
"""
import sys

def main():
    print("=" * 50)
    print(f"项目: {description}")
    print("=" * 50)
    print("这是一个示例项目框架。")
    print("请根据需求修改 main.py 文件。")

if __name__ == "__main__":
    main()
'''

    def _general_html(self, description: str) -> str:
        return f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{description}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: -apple-system, sans-serif; background: #f5f5f5; padding: 20px; }}
        .container {{ max-width: 800px; margin: 0 auto; background: #fff; border-radius: 12px; padding: 30px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }}
        h1 {{ color: #333; margin-bottom: 20px; }}
        p {{ color: #666; line-height: 1.8; }}
    </style>
</head>
<body>
    <div class="container">
        <h1>{description}</h1>
        <p>这是一个自动生成的网页项目。</p>
    </div>
</body>
</html>'''

    def _general_js(self, description: str) -> str:
        return f'''/**
 * {description}
 */
console.log("{description}");
console.log("项目已创建，请根据需求修改。");
'''


class ListDirTool(Tool):
    """列出目录工具"""
    
    def __init__(self):
        super().__init__("list_dir", "List files and directories at a given path. Input: a directory path, e.g. 'C:\\Users', '.', '/home'. Defaults to current directory if empty.")
    
    def run(self, input: str) -> str:
        import os
        try:
            path = input.strip() or "."
            items = os.listdir(path)
            result = []
            for item in items[:50]:
                full_path = os.path.join(path, item)
                if os.path.isdir(full_path):
                    result.append(f"[DIR] {item}/")
                else:
                    size = os.path.getsize(full_path)
                    if size > 1024*1024:
                        result.append(f"[FILE] {item} ({size//1024//1024}MB)")
                    elif size > 1024:
                        result.append(f"[FILE] {item} ({size//1024}KB)")
                    else:
                        result.append(f"[FILE] {item} ({size}B)")
            return "\n".join(result) if result else "目录为空"
        except Exception as e:
            return f"读取失败: {e}"


class WebSearchTool(Tool):
    """搜索工具（DuckDuckGo）"""
    
    def __init__(self):
        super().__init__("web_search", "Search the internet for information. Input: search query string. Returns top 5 results with titles, snippets, and URLs.")
    
    def run(self, input: str) -> str:
        try:
            from duckduckgo_search import DDGS
            with DDGS() as ddgs:
                results = list(ddgs.text(input, max_results=5))
            if not results:
                return f"未找到关于 '{input}' 的搜索结果"
            output = f"搜索 '{input}' 的结果:\n\n"
            for i, r in enumerate(results, 1):
                title = r.get("title", "")
                body = r.get("body", "")
                url = r.get("href", "")
                output += f"{i}. {title}\n   {body}\n   链接: {url}\n\n"
            return output.strip()
        except Exception as e:
            return f"搜索出错: {e}"


class DateTimeTool(Tool):
    """日期时间工具"""
    
    def __init__(self):
        super().__init__("datetime", "Get the current date and time. Input: ignored, any string is fine.")
    
    def run(self, input: str = "") -> str:
        from datetime import datetime
        return f"当前时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

ALL_TOOLS = [
    CalculatorTool(), OpenAppTool(), OpenUrlTool(), OpenFileTool(),
    RunCommandTool(), ListDirTool(), WebSearchTool(), DateTimeTool(),
    CodeGeneratorTool()
]

TOOL_MAP = {t.name: t for t in ALL_TOOLS}

def execute_tool_batch(tool_calls: list) -> list:
    results = []
    for call in tool_calls:
        tool_name = call.get("tool_name", "")
        tool_input = call.get("input", "")
        if tool_name in TOOL_MAP:
            try:
                result = TOOL_MAP[tool_name].run(tool_input)
                results.append({"tool": tool_name, "result": result, "success": True})
            except Exception as e:
                results.append({"tool": tool_name, "result": str(e), "success": False})
        else:
            results.append({"tool": tool_name, "result": f"工具 '{tool_name}' 不存在", "success": False})
    return results

class SimpleAgent:
    """简单 ReAct 智能体"""
    
    def __init__(self, llm: MultiProviderLLM, tools: List[Tool] = None, system_prompt: str = ""):
        self.llm = llm
        self.tools = tools or [CalculatorTool(), DateTimeTool(), WebSearchTool()]
        self.tool_map = {t.name: t for t in self.tools}
        self.system_prompt = system_prompt or self._default_prompt()
    
    def _default_prompt(self) -> str:
        tool_desc = "\n".join([f"- {t.name}: {t.description}" for t in self.tools])
        return f"""你是一个智能助手，可以使用以下工具:
{tool_desc}

使用工具时，请按以下格式输出:
Thought: 我需要...
Action: 工具名
Action Input: 输入内容

不使用工具时，直接回答用户问题。
请用中文回答。"""
    
    def run(self, user_input: str) -> str:
        messages = [
            SystemMessage(content=self.system_prompt),
            HumanMessage(content=user_input)
        ]
        
        response = self.llm._generate(messages)
        
        # 简单解析 Action
        if "Action:" in response:
            try:
                lines = response.split("\n")
                action = None
                action_input = None
                for line in lines:
                    if line.strip().startswith("Action:"):
                        action = line.split("Action:")[-1].strip()
                    elif line.strip().startswith("Action Input:"):
                        action_input = line.split("Action Input:")[-1].strip()
                
                if action and action in self.tool_map:
                    tool_result = self.tool_map[action].run(action_input or "")
                    messages.append(AIMessage(content=response))
                    messages.append(HumanMessage(content=f"工具返回结果: {tool_result}\n请基于此结果回答用户问题。"))
                    final_response = self.llm._generate(messages)
                    return final_response
            except Exception:
                pass
        
        return response


# ============ 5. RAG（检索增强生成） ============

class RAGEngine:
    """RAG 引擎"""
    
    def __init__(self, collection_dir: str = "./chroma_db"):
        import threading
        self.collection_dir = collection_dir
        self.embeddings = None
        self.vectorstore = None
        self._cleared = False
        self._lock = threading.Lock()
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=500, chunk_overlap=50, separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?"]
        )
    
    def _get_embeddings(self):
        if self.embeddings is None:
            self.embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/all-MiniLM-L6-v2",
                model_kwargs={"device": "cpu"}
            )
        return self.embeddings
    
    def add_text(self, text: str, metadata: dict = None) -> int:
        """添加文本到向量库"""
        self._cleared = False
        self._ensure_vectorstore()
        docs = self.text_splitter.create_documents([text], metadatas=[metadata or {}])
        if self.vectorstore is None:
            self.vectorstore = Chroma.from_documents(docs, self._get_embeddings(), persist_directory=self.collection_dir)
        else:
            self.vectorstore.add_documents(docs)
        return len(docs)
    
    def add_file(self, file_path: str) -> int:
        """加载文件并添加到向量库"""
        ext = os.path.splitext(file_path)[1].lower()
        loaders = {
            ".txt": lambda p: TextLoader(p, encoding="utf-8"),
            ".pdf": lambda p: PyPDFLoader(p),
            ".csv": lambda p: CSVLoader(p),
            ".md": lambda p: TextLoader(p, encoding="utf-8"),
            ".html": lambda p: TextLoader(p, encoding="utf-8"),
            ".htm": lambda p: TextLoader(p, encoding="utf-8"),
            ".json": lambda p: TextLoader(p, encoding="utf-8"),
            ".xml": lambda p: TextLoader(p, encoding="utf-8"),
            ".log": lambda p: TextLoader(p, encoding="utf-8"),
            ".docx": self._load_docx,
            ".xlsx": self._load_xlsx,
            ".xls": self._load_xlsx,
            ".pptx": self._load_pptx,
            ".ppt": self._load_pptx,
        }

        if ext not in loaders:
            try:
                loader = TextLoader(file_path, encoding="utf-8")
                docs = loader.load()
            except Exception:
                with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                    from langchain_core.documents import Document
                    docs = [Document(page_content=f.read(), metadata={"source": file_path})]
        else:
            loader = loaders[ext](file_path)
            docs = loader.load() if callable(getattr(loader, "load", None)) else loader

        chunks = self.text_splitter.split_documents(docs)

        with self._lock:
            if self._cleared:
                self._cleared = False
                self._ensure_vectorstore()
            if self.vectorstore is None:
                self.vectorstore = Chroma.from_documents(chunks, self._get_embeddings(), persist_directory=self.collection_dir)
            else:
                self.vectorstore.add_documents(chunks)

        return len(chunks)

    def add_files(self, file_paths: list) -> int:
        """批量添加文件"""
        total = 0
        for fp in file_paths:
            total += self.add_file(fp)
        return total

    def _load_docx(self, path):
        from langchain_core.documents import Document
        from docx import Document as DocxDocument
        doc = DocxDocument(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return [Document(page_content="\n".join(paragraphs), metadata={"source": path})]

    def _load_xlsx(self, path):
        from langchain_core.documents import Document
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                line = " | ".join(str(c) if c is not None else "" for c in row)
                if line.strip():
                    texts.append(line)
        wb.close()
        return [Document(page_content="\n".join(texts), metadata={"source": path})]

    def _load_pptx(self, path):
        from langchain_core.documents import Document
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return [Document(page_content="\n".join(texts), metadata={"source": path})]
    
    def _ensure_vectorstore(self):
        if self._cleared:
            return
        if self.vectorstore is None and os.path.exists(self.collection_dir) and os.listdir(self.collection_dir):
            try:
                self.vectorstore = Chroma(persist_directory=self.collection_dir, embedding_function=self._get_embeddings())
            except Exception:
                pass
    
    def search(self, query: str, top_k: int = 3) -> List[Dict]:
        """检索相关文档"""
        if self._cleared:
            return []
        with self._lock:
            self._ensure_vectorstore()
            if self.vectorstore is None:
                return []
            results = self.vectorstore.similarity_search_with_score(query, k=top_k)
        return [{"content": doc.page_content, "score": float(score), "metadata": doc.metadata} for doc, score in results]
    
    def clear(self) -> bool:
        """清空知识库所有数据"""
        import chromadb, gc
        self._cleared = True
        with self._lock:
            self.vectorstore = None
            self.embeddings = None
            if os.path.exists(self.collection_dir):
                try:
                    client = chromadb.PersistentClient(path=self.collection_dir)
                    for col in client.list_collections():
                        client.delete_collection(col.name)
                    del client
                    gc.collect()
                except Exception:
                    pass
                import shutil, time
                for attempt in range(5):
                    try:
                        shutil.rmtree(self.collection_dir)
                        break
                    except (PermissionError, OSError):
                        time.sleep(1)
                        gc.collect()
                os.makedirs(self.collection_dir, exist_ok=True)
        return True
    
    def rag_query(self, query: str, llm=None, top_k: int = 3) -> str:
        """RAG 查询：检索 + 生成"""
        results = self.search(query, top_k)
        if not results:
            return "知识库中暂无相关内容。"
        
        context = "\n\n".join([r["content"] for r in results])
        prompt = f"""基于以下参考资料回答问题。如果资料中没有相关信息，请说明。

参考资料:
{context}

问题: {query}
回答:"""
        
        if llm is None:
            try:
                from llm_config import get_llm_router
                llm = get_llm_router()
            except Exception:
                pass
        
        if llm is None:
            return "LLM 未配置，无法生成回答。"
        
        provider_name = getattr(llm, 'default_provider', None)
        
        import asyncio, threading
        result_container = [None]
        error_container = [None]

        def _run():
            new_loop = asyncio.new_event_loop()
            asyncio.set_event_loop(new_loop)
            try:
                async def _call():
                    r = ""
                    async for chunk in llm.chat_stream([{"role": "user", "content": prompt}], provider=provider_name, enable_thinking=False):
                        if chunk == "[DONE]" or (isinstance(chunk, str) and chunk.startswith('{"type"')):
                            continue
                        r += chunk
                    return r
                result_container[0] = new_loop.run_until_complete(_call())
            except Exception as e:
                error_container[0] = str(e)
            finally:
                new_loop.close()

        t = threading.Thread(target=_run)
        t.start()
        t.join(timeout=60)

        if error_container[0]:
            return f"AI 生成失败: {error_container[0]}"
        return result_container[0] or "AI 未返回内容。"
    
    def get_stats(self) -> Dict:
        if self.vectorstore is not None:
            try:
                return {"total_chunks": self.vectorstore._collection.count()}
            except Exception:
                pass
        if os.path.exists(self.collection_dir) and os.listdir(self.collection_dir):
            try:
                client = chromadb.PersistentClient(path=self.collection_dir)
                collections = client.list_collections()
                total = sum(c.count() for c in collections)
                return {"total_chunks": total}
            except Exception:
                pass
        return {"total_chunks": 0}


# ============ 6. 文档处理工具 ============

class DocumentProcessor:
    """文档处理器"""
    
    @staticmethod
    def split_text(text: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size, chunk_overlap=chunk_overlap
        )
        return splitter.split_text(text)
    
    @staticmethod
    def load_and_split(file_path: str, chunk_size: int = 500) -> List[Dict]:
        ext = os.path.splitext(file_path)[1].lower()
        loaders = {
            ".txt": lambda p: TextLoader(p, encoding="utf-8"),
            ".pdf": lambda p: PyPDFLoader(p),
            ".csv": lambda p: CSVLoader(p),
        }
        if ext not in loaders:
            raise ValueError(f"不支持: {ext}")
        
        docs = loaders[ext](file_path).load()
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=50)
        chunks = splitter.split_documents(docs)
        return [{"content": c.page_content, "metadata": c.metadata} for c in chunks]


# ============ 全局实例 ============

_rag_engine = None

def get_rag_engine() -> RAGEngine:
    global _rag_engine
    if _rag_engine is None:
        _rag_engine = RAGEngine()
    return _rag_engine

def get_llm(provider: str = None, **kwargs) -> MultiProviderLLM:
    return MultiProviderLLM(provider=provider, **kwargs)
