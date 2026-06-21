from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Query
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional
from fastapi.responses import StreamingResponse
from datetime import datetime
import uvicorn
import os
import json
import uuid
import threading
import queue

from orchestrator import Orchestrator

app = FastAPI(title="AI应用通用框架", version="1.0.0")

@app.on_event("startup")
async def preload_embeddings():
    def _load():
        try:
            from langchain_engine import get_rag_engine
            rag = get_rag_engine()
            rag._get_embeddings()
            print("[STARTUP] Embedding model loaded")
        except Exception as e:
            print(f"[STARTUP] Embedding preload failed: {e}")
    threading.Thread(target=_load, daemon=True).start()

@app.exception_handler(Exception)
async def global_exception_handler(request, exc):
    import traceback
    traceback.print_exc()
    from fastapi.responses import JSONResponse
    return JSONResponse(status_code=500, content={"detail": str(exc)})
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

ENV_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")

# ============ MongoDB ============

from pymongo import MongoClient
MONGO_URL = os.environ.get("MONGODB_URL", "mongodb://ai_mongo:ai_mongo_2024@localhost:27017")
mongo_client = MongoClient(MONGO_URL, serverSelectionTimeoutMS=3000)
db = mongo_client["ai_framework"]
conversations_col = db["conversations"]
agents_col = db["agents"]
settings_col = db["settings"]
tasks_col = db["tasks"]

# ============ Orchestrator 初始化 ==============

_orchestrator = None

def get_orchestrator():
    global _orchestrator
    if _orchestrator is None:
        from langchain_engine import ALL_TOOLS
        tool_map = {t.name: t for t in ALL_TOOLS}
        llm = get_llm_router()
        _orchestrator = Orchestrator(llm, tool_map, db)
        print(f"[Orchestrator] Created with providers: {list(llm.providers.keys())}, default: {llm.default_provider}")
    return _orchestrator

# ============ 数据库配置读写 ============

def get_settings_from_db() -> dict:
    try:
        doc = settings_col.find_one({"_id": "app_settings"})
        if doc:
            doc.pop("_id", None)
            return doc
    except Exception:
        pass
    return {}

def save_settings_to_db(updates: dict):
    try:
        settings_col.update_one(
            {"_id": "app_settings"},
            {"$set": updates},
            upsert=True
        )
    except Exception as e:
        print(f"Warning: Could not save settings to DB: {e}")

def get_provider_config(name: str) -> dict:
    """从数据库读取提供商配置"""
    s = get_settings_from_db()
    key = s.get(f"{name}_key", "")
    url = s.get(f"{name}_url", "")
    model = s.get(f"{name}_model", "")
    return {"api_key": key, "base_url": url, "model": model}

def get_default_provider() -> str:
    s = get_settings_from_db()
    return s.get("default_provider", "agnes")

# ============ 动态提供商管理 ============

def get_llm_router():
    from llm_config import get_llm_router as _get
    return _get()

def register_provider_dynamic(name: str, api_key: str, base_url: str, model: str = None):
    from llm_providers import DeepSeekProvider, MiMoProvider, OpenAIProvider, QwenProvider, GLMProvider, AgnesProvider
    router = get_llm_router()
    provider_map = {"deepseek": DeepSeekProvider, "mimo": MiMoProvider, "openai": OpenAIProvider, "qwen": QwenProvider, "glm": GLMProvider, "agnes": AgnesProvider, "custom": OpenAIProvider}
    if name in provider_map and api_key and not api_key.startswith("your-"):
        kwargs = {"api_key": api_key, "base_url": base_url}
        if model: kwargs["model"] = model
        router.providers[name] = provider_map[name](**kwargs)
        return True
    return False

# ============ 数据模型 ============

class ChatRequest(BaseModel):
    message: str
    context: Optional[dict] = None

class ChatWithProviderRequest(BaseModel):
    message: str
    provider: str = None
    model: str = None
    stream: bool = False
    context: Optional[dict] = None

class ChatStreamRequest(BaseModel):
    conversation_id: Optional[str] = None
    message: str
    provider: Optional[str] = None
    model: Optional[str] = None
    system_prompt: Optional[str] = None
    images: Optional[List[str]] = None
    enable_thinking: Optional[bool] = True

class AgentCreate(BaseModel):
    name: str
    desc: Optional[str] = ""
    icon: Optional[str] = "🤖"
    bgColor: Optional[str] = "#e8f0fe"
    category: Optional[str] = "work"
    systemPrompt: str

class ChatResponse(BaseModel):
    response: str
    confidence: float
    sources: Optional[List[str]] = None

class ConversationCreate(BaseModel):
    title: Optional[str] = "新对话"
    provider: Optional[str] = None

class ContentRequest(BaseModel):
    topic: str
    style: Optional[str] = "专业分析"
    length: Optional[str] = "medium"

class ContentResponse(BaseModel):
    content: str
    word_count: int
    reading_time: str

class DocumentRequest(BaseModel):
    topic: str
    format: Optional[str] = "ppt"
    style: Optional[str] = "商务"
    pages: Optional[int] = 10
    extra: Optional[str] = ""

class DocumentResponse(BaseModel):
    status: str
    filename: str
    message: Optional[str] = None

class RecommendationRequest(BaseModel):
    user_id: str
    context: Optional[dict] = None
    top_k: Optional[int] = 10

class RecommendationResponse(BaseModel):
    recommendations: List[dict]
    scores: List[float]

class SettingsRequest(BaseModel):
    deepseek_key: Optional[str] = None
    deepseek_url: Optional[str] = None
    deepseek_model: Optional[str] = None
    mimo_key: Optional[str] = None
    mimo_url: Optional[str] = None
    mimo_model: Optional[str] = None
    qwen_key: Optional[str] = None
    qwen_url: Optional[str] = None
    qwen_model: Optional[str] = None
    glm_key: Optional[str] = None
    glm_url: Optional[str] = None
    glm_model: Optional[str] = None
    agnes_key: Optional[str] = None
    agnes_url: Optional[str] = None
    agnes_model: Optional[str] = None
    openai_key: Optional[str] = None
    openai_url: Optional[str] = None
    openai_model: Optional[str] = None
    custom_key: Optional[str] = None
    custom_url: Optional[str] = None
    custom_model: Optional[str] = None
    default_provider: Optional[str] = None
    temperature: Optional[float] = None
    max_tokens: Optional[int] = None
    # 图像生成配置
    image_provider: Optional[str] = None
    agnes_image_model_1: Optional[str] = None
    agnes_image_model_2: Optional[str] = None
    custom_image_key: Optional[str] = None
    custom_image_url: Optional[str] = None
    custom_image_model: Optional[str] = None
    image_size: Optional[str] = None
    image_num: Optional[int] = None
    image_quality: Optional[str] = None
    image_negative_prompt: Optional[str] = None
    # 图像生成高级参数
    image_style: Optional[str] = None
    image_seed: Optional[int] = None
    image_guidance_scale: Optional[float] = None
    # 视频生成配置
    video_provider: Optional[str] = None
    agnes_video_model: Optional[str] = None
    custom_video_key: Optional[str] = None
    custom_video_url: Optional[str] = None
    custom_video_model: Optional[str] = None
    video_duration: Optional[int] = None
    video_fps: Optional[int] = None
    video_resolution: Optional[str] = None
    video_negative_prompt: Optional[str] = None
    video_seed: Optional[int] = None

# ============ 对话管理 API ============

@app.get("/api/v1/conversations")
async def list_conversations():
    convs = list(conversations_col.find({}, {"messages": 0}).sort("updated_at", -1))
    for c in convs:
        c["_id"] = str(c["_id"])
    return {"conversations": convs}

@app.post("/api/v1/conversations")
async def create_conversation(req: ConversationCreate):
    conv_id = str(uuid.uuid4())[:8]
    doc = {
        "_id": conv_id,
        "title": req.title,
        "provider": req.provider,
        "messages": [],
        "created_at": datetime.now().isoformat(),
        "updated_at": datetime.now().isoformat()
    }
    conversations_col.insert_one(doc)
    return {"id": conv_id, "title": doc["title"], "provider": doc["provider"]}

@app.get("/api/v1/conversations/{conv_id}")
async def get_conversation(conv_id: str):
    conv = conversations_col.find_one({"_id": conv_id})
    if not conv:
        raise HTTPException(status_code=404, detail="对话不存在")
    conv["_id"] = str(conv["_id"])
    return conv

@app.delete("/api/v1/conversations/{conv_id}")
async def delete_conversation(conv_id: str):
    conversations_col.delete_one({"_id": conv_id})
    return {"status": "ok"}

# ============ 智能体 API ============

@app.get("/api/v1/agents")
async def list_agents():
    agents = list(agents_col.find({}))
    for a in agents:
        a["_id"] = str(a["_id"])
    return {"agents": agents}

@app.post("/api/v1/agents")
async def create_agent(req: AgentCreate):
    agent_id = str(uuid.uuid4())[:8]
    doc = {
        "_id": agent_id,
        "name": req.name,
        "desc": req.desc,
        "icon": req.icon,
        "bgColor": req.bgColor,
        "category": req.category,
        "systemPrompt": req.systemPrompt,
        "chats": 0,
        "author": "管理员",
        "created_at": datetime.now().isoformat()
    }
    agents_col.insert_one(doc)
    doc["_id"] = agent_id
    return doc

@app.delete("/api/v1/agents/{agent_id}")
async def delete_agent(agent_id: str):
    agents_col.delete_one({"_id": agent_id})
    return {"status": "ok"}

# ============ 工具定义 ============

def get_tool_descriptions():
    from langchain_engine import CalculatorTool, DateTimeTool, OpenUrlTool, RunCommandTool, WebSearchTool, CodeGeneratorTool
    tools = [CalculatorTool(), DateTimeTool(), OpenUrlTool(), RunCommandTool(), WebSearchTool(), CodeGeneratorTool()]
    return tools

def get_tool_map():
    return {t.name: t for t in get_tool_descriptions()}

def get_tool_prompt_suffix():
    tools = get_tool_descriptions()
    tool_lines = "\n".join([f"- {t.name}: {t.description}" for t in tools])
    return f"""

你也可以使用以下工具来帮助用户:
{tool_lines}

当你需要使用工具时，请在回复中按以下格式输出（仅限确实需要工具时）:
Thought: 我需要...
Action: 工具名
Action Input: 输入内容

重要规则:
- 只使用上面列出的工具，不要编造不存在的工具
- 不要滥用工具。只有当用户明确要求或问题确实需要工具才能回答时才使用
- code_generate 工具的输入格式是: 项目名|简短描述|语言，不要把整个代码放在输入中
- open_url 工具的输入格式是: URL地址
- calculator 工具的输入格式是: 数学表达式
- run_command 工具的输入格式是: 命令字符串"""

def parse_tool_call(text: str):
    import re
    
    # 格式1: <tool_call>name<tool_call> + <tool_input>input</tool_input>
    name_match = re.search(r'<tool_call>(\w+)<tool_call>', text)
    input_match = re.search(r'<tool_input>(.*?)</tool_input>', text, re.DOTALL)
    if name_match and input_match:
        thinking_match = re.search(r'<thinking>(.*?)</thinking>', text, re.DOTALL)
        return {
            "thinking": thinking_match.group(1).strip() if thinking_match else "",
            "tool": name_match.group(1).strip(),
            "input": input_match.group(1).strip()
        }
    
    # 格式2: Action: tool_name + Action Input: input (ReAct格式)
    # 用逐行解析，兼容多行输入和转义字符
    lines = text.split('\n')
    action_name = None
    action_input_lines = []
    capture_input = False
    
    for line in lines:
        stripped = line.strip()
        if stripped.startswith('Action:'):
            action_name = stripped.replace('Action:', '').strip()
        elif stripped.startswith('Action Input:'):
            capture_input = True
            # Action Input: 后面的内容在同一行
            rest = stripped.replace('Action Input:', '').strip()
            if rest:
                action_input_lines.append(rest)
        elif capture_input:
            # 如果下一行不是新的标记，继续收集
            if stripped and not stripped.startswith(('Action:', 'Thought:', 'Observation:', 'Final Answer:')):
                action_input_lines.append(stripped)
            else:
                capture_input = False
    
    if action_name and action_input_lines:
        thinking_match = re.search(r'Thought:\s*(.*?)(?=\n|$)', text, re.DOTALL)
        return {
            "thinking": thinking_match.group(1).strip() if thinking_match else "",
            "tool": action_name,
            "input": ' '.join(action_input_lines)
        }
    
    return None

# ============ 流式对话 ============

@app.post("/api/v1/chat/stream")
async def chat_stream(request: ChatStreamRequest):
    conv_id = request.conversation_id

    history = []
    if conv_id:
        conv = conversations_col.find_one({"_id": conv_id})
        if conv:
            history = conv.get("messages", [])

    if conv_id:
        conversations_col.update_one(
            {"_id": conv_id},
            {"$push": {"messages": {"role": "user", "content": request.message}},
             "$set": {"updated_at": datetime.now().isoformat(), "provider": request.provider}}
        )

    try:
        from recommendation_engine import add_interaction
        add_interaction("default_user", request.message)
    except Exception:
        pass

    collected = [""]

    async def generate():
        try:
            router = get_llm_router()
            
            messages = []
            if request.system_prompt:
                messages.append({"role": "system", "content": request.system_prompt})
            
            for msg in history:
                messages.append({"role": msg.get("role", "user"), "content": msg.get("content", "")})
            
            if request.images:
                content = [{"type": "text", "text": request.message}]
                for img in request.images:
                    content.append({"type": "image_url", "image_url": {"url": img}})
                messages.append({"role": "user", "content": content})
            else:
                messages.append({"role": "user", "content": request.message})
            
            async for chunk in router.chat_stream(messages, provider=request.provider, model=request.model, enable_thinking=request.enable_thinking):
                if chunk == "[DONE]":
                    continue
                if chunk.startswith('{"type": "reasoning"'):
                    try:
                        r = json.loads(chunk)
                        yield f"data: {json.dumps({'reasoning': r['content']})}\n\n"
                    except:
                        pass
                else:
                    collected[0] += chunk
                    yield f"data: {json.dumps({'content': chunk})}\n\n"

            if collected[0] and len(collected[0]) > 5:
                try:
                    suggest_prompt = [
                        {"role": "system", "content": "你是一个助手。根据用户的问题和AI的回答，生成3个相关的后续问题，让用户可以继续深入了解。只输出问题，每行一个，不要编号，不要其他内容。"},
                        {"role": "user", "content": f"用户问题：{request.message}\nAI回答：{collected[0][:500]}"}
                    ]
                    suggest_response = await router.chat(suggest_prompt, provider=request.provider, model=request.model)
                    suggestions = [s.strip() for s in suggest_response.content.strip().split("\n") if s.strip() and len(s.strip()) > 3]
                    suggestions = suggestions[:3]
                    if suggestions:
                        yield f"data: {json.dumps({'suggestions': suggestions})}\n\n"
                except Exception:
                    pass

        except Exception:
            pass

        # 保存对话
        if conv_id and collected[0]:
            try:
                conversations_col.update_one(
                    {"_id": conv_id},
                    {"$push": {"messages": {"role": "assistant", "content": collected[0]}},
                     "$set": {"updated_at": datetime.now().isoformat()}}
                )
            except Exception:
                pass

        yield "data: [DONE]\n\n"

    async def save_later():
        await asyncio.sleep(5)
        if conv_id and collected[0]:
            try:
                conv = conversations_col.find_one({"_id": conv_id})
                if conv:
                    msgs = conv.get("messages", [])
                    has_assistant = any(m.get("role") == "assistant" for m in msgs)
                    if not has_assistant:
                        conversations_col.update_one(
                            {"_id": conv_id},
                            {"$push": {"messages": {"role": "assistant", "content": collected[0]}},
                             "$set": {"updated_at": datetime.now().isoformat()}}
                        )
            except Exception:
                pass

    import asyncio
    background_task = asyncio.get_event_loop().create_task(save_later())

    return StreamingResponse(generate(), media_type="text/event-stream", headers={
        "Cache-Control": "no-cache", "Connection": "keep-alive", "X-Accel-Buffering": "no"
    })

# ============ 任务管理 API ============

@app.get("/api/v1/tasks/{task_id}")
async def get_task(task_id: str):
    orchestrator = get_orchestrator()
    if not orchestrator.task_manager:
        raise HTTPException(status_code=500, detail="任务管理器未初始化")
    task = await orchestrator.task_manager.get_task(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="任务不存在")
    return task.model_dump() if hasattr(task, 'model_dump') else task.dict()

@app.get("/api/v1/tasks")
async def list_tasks(conversation_id: str = None):
    orchestrator = get_orchestrator()
    if not orchestrator.task_manager:
        raise HTTPException(status_code=500, detail="任务管理器未初始化")
    if conversation_id:
        tasks = await orchestrator.task_manager.get_pending_tasks(conversation_id)
    else:
        cursor = orchestrator.task_manager.collection.find({}).sort("created_at", -1).limit(50)
        tasks = []
        for doc in cursor:
            doc.pop("_id", None)
            tasks.append(doc)
    return {"tasks": tasks}

@app.post("/api/v1/tasks/{task_id}/resume")
async def resume_task(task_id: str):
    orchestrator = get_orchestrator()

    async def generate():
        async for event in orchestrator.resume_task(task_id):
            yield f"data: {json.dumps(event)}\n\n"
        yield "data: [DONE]\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream", headers={
        "Cache-Control": "no-cache", "Connection": "keep-alive", "X-Accel-Buffering": "no"
    })

# ============ 原有 API ============

@app.get("/")
async def root():
    return {"message": "AI应用通用框架", "version": "1.0.0", "docs": "/docs"}

@app.get("/api/v1/health")
async def health_check():
    return {"status": "healthy", "services": {"api": "running"}}

@app.get("/api/v1/dashboard")
async def get_dashboard():
    from datetime import datetime, timedelta
    now = datetime.now()
    today_start = now.replace(hour=0, minute=0, second=0, microsecond=0)

    stats = {}
    try:
        stats["conversations_today"] = conversations_col.count_documents({"created_at": {"$gte": today_start.isoformat()}})
        stats["conversations_total"] = conversations_col.count_documents({})
    except Exception:
        stats["conversations_today"] = 0
        stats["conversations_total"] = 0

    try:
        from workflow_engine import _get_cols
        wcol, rcol = _get_cols()
        stats["workflow_runs_today"] = rcol.count_documents({"created_at": {"$gte": today_start.isoformat()}})
        stats["workflow_runs_total"] = rcol.count_documents({})
    except Exception:
        stats["workflow_runs_today"] = 0
        stats["workflow_runs_total"] = 0

    try:
        from workflow_engine import _get_cols
        wcol, _ = _get_cols()
        stats["workflows_total"] = wcol.count_documents({})
    except Exception:
        stats["workflows_total"] = 0

    try:
        stats["agents_total"] = agents_col.count_documents({})
    except Exception:
        stats["agents_total"] = 0

    recent_chats = []
    try:
        for doc in conversations_col.find().sort("created_at", -1).limit(10):
            msgs = doc.get("messages", [])
            last_msg = msgs[-1].get("content", "") if msgs else ""
            recent_chats.append({
                "time": doc.get("created_at", "")[:16].replace("T", " "),
                "model": doc.get("provider", "unknown"),
                "message": last_msg[:80],
                "title": doc.get("title", "")
            })
    except Exception:
        pass

    system_status = []
    try:
        mongo_client.admin.command("ping")
        system_status.append({"name": "MongoDB", "ok": True})
    except Exception:
        system_status.append({"name": "MongoDB", "ok": False})

    try:
        import httpx
        settings_doc = settings_col.find_one({"_id": "app_settings"}) or {}
        providers = []
        provider_map = {
            "deepseek": ("deepseek_key", "deepseek_url"),
            "mimo": ("mimo_key", "mimo_url"),
            "openai": ("openai_key", "openai_url"),
            "qwen": ("qwen_key", "qwen_url"),
            "glm": ("glm_key", "glm_url"),
            "agnes": ("agnes_key", "agnes_url"),
        }
        for name, (key_field, url_field) in provider_map.items():
            api_key = settings_doc.get(key_field, "")
            base_url = settings_doc.get(url_field, "")
            available = bool(api_key)
            providers.append({"name": name, "available": available})
        system_status.append({"name": "LLM Providers", "ok": any(p["available"] for p in providers)})
    except Exception:
        providers = [{"name": n, "available": False} for n in ["deepseek", "mimo", "openai", "qwen", "glm", "agnes"]]
        system_status.append({"name": "LLM Providers", "ok": False})

    try:
        today_chats = stats["conversations_today"]
        today_runs = stats["workflow_runs_today"]
        total_today = today_chats + today_runs
        usage = [
            {"name": "AI 对话", "value": today_chats, "pct": min(100, today_chats * 5) if total_today > 0 else 0},
            {"name": "工作流", "value": today_runs, "pct": min(100, today_runs * 10) if total_today > 0 else 0},
        ]
    except Exception:
        usage = []

    return {
        "stats": stats,
        "recent_chats": recent_chats,
        "system_status": system_status,
        "providers": providers,
        "usage": usage
    }

@app.post("/api/v1/settings")
async def save_settings(request: SettingsRequest):
    updates = {}
    mapping = {
        "deepseek_key": "deepseek_key", "deepseek_url": "deepseek_url", "deepseek_model": "deepseek_model",
        "mimo_key": "mimo_key", "mimo_url": "mimo_url", "mimo_model": "mimo_model",
        "qwen_key": "qwen_key", "qwen_url": "qwen_url", "qwen_model": "qwen_model",
        "glm_key": "glm_key", "glm_url": "glm_url", "glm_model": "glm_model",
        "agnes_key": "agnes_key", "agnes_url": "agnes_url", "agnes_model": "agnes_model",
        "openai_key": "openai_key", "openai_url": "openai_url", "openai_model": "openai_model",
        "custom_key": "custom_key", "custom_url": "custom_url", "custom_model": "custom_model",
        "default_provider": "default_provider", "temperature": "temperature", "max_tokens": "max_tokens",
        "image_provider": "image_provider",
        "agnes_image_model_1": "agnes_image_model_1", "agnes_image_model_2": "agnes_image_model_2",
        "custom_image_key": "custom_image_key", "custom_image_url": "custom_image_url", "custom_image_model": "custom_image_model",
        "image_size": "image_size", "image_num": "image_num", "image_quality": "image_quality",
        "image_negative_prompt": "image_negative_prompt",
        "image_style": "image_style", "image_seed": "image_seed", "image_guidance_scale": "image_guidance_scale",
        "video_provider": "video_provider",
        "agnes_video_model": "agnes_video_model",
        "custom_video_key": "custom_video_key", "custom_video_url": "custom_video_url", "custom_video_model": "custom_video_model",
        "video_duration": "video_duration", "video_fps": "video_fps", "video_resolution": "video_resolution",
        "video_negative_prompt": "video_negative_prompt", "video_seed": "video_seed",
    }
    for field, db_key in mapping.items():
        value = getattr(request, field, None)
        if value is not None:
            updates[db_key] = str(value)
    
    if updates:
        save_settings_to_db(updates)
    
    # 刷新 LLM 路由器
    from llm_config import refresh_llm_router
    refresh_llm_router()
    
    return {"status": "ok", "message": "配置已保存并生效"}

@app.get("/api/v1/settings")
async def get_settings():
    db_settings = get_settings_from_db()
    defaults = {
        "deepseek_key": "", "deepseek_url": "https://api.deepseek.com", "deepseek_model": "deepseek-chat",
        "mimo_key": "", "mimo_url": "https://api.mimo.xiaomi.com", "mimo_model": "mimo-auto",
        "qwen_key": "", "qwen_url": "https://dashscope.aliyuncs.com", "qwen_model": "qwen-turbo",
        "glm_key": "", "glm_url": "https://open.bigmodel.cn/api/paas/v4", "glm_model": "glm-4-flash",
        "agnes_key": "", "agnes_url": "https://apihub.agnes-ai.com", "agnes_model": "agnes-2.0-flash",
        "openai_key": "", "openai_url": "https://api.openai.com", "openai_model": "gpt-3.5-turbo",
        "custom_key": "", "custom_url": "", "custom_model": "",
        "default_provider": "agnes",
        "image_provider": "agnes",
        "agnes_image_model_1": "agnes-image-2.0-flash", "agnes_image_model_2": "agnes-image-2.1-flash",
        "custom_image_key": "", "custom_image_url": "", "custom_image_model": "",
        "image_size": "original", "image_num": 1, "image_quality": "fhd", "image_negative_prompt": "",
        "image_style": "", "image_seed": None, "image_guidance_scale": None,
        "video_provider": "agnes",
        "agnes_video_model": "agnes-video-v2.0",
        "custom_video_key": "", "custom_video_url": "", "custom_video_model": "",
        "video_duration": 5, "video_fps": 24, "video_resolution": "720p",
        "video_negative_prompt": "", "video_seed": None,
    }
    result = {**defaults, **db_settings}
    result.pop("_id", None)
    # 类型修正
    try: result["temperature"] = float(result.get("temperature", 0.7))
    except: result["temperature"] = 0.7
    try: result["max_tokens"] = int(result.get("max_tokens", 2048))
    except: result["max_tokens"] = 2048
    return result

@app.get("/api/v1/llm/models")
async def get_model_list(api_key: str, base_url: str):
    """获取模型列表"""
    import httpx
    try:
        url = base_url.rstrip('/')
        if '/api/paas/v4' in url:
            endpoint = f"{url}/models"
        elif url.endswith('/v1') or url.endswith('/v4') or '/v1/' in url or '/v4/' in url:
            endpoint = f"{url}/models"
        else:
            endpoint = f"{url}/v1/models"
        
        async with httpx.AsyncClient(timeout=15.0) as client:
            response = await client.get(
                endpoint,
                headers={"Authorization": f"Bearer {api_key}"}
            )
            if response.status_code == 200:
                data = response.json()
                models = [m.get("id", "") for m in data.get("data", []) if m.get("id")]
                return {"status": "ok", "models": models}
            else:
                return {"status": "error", "message": f"HTTP {response.status_code}: {response.text[:200]}"}
    except httpx.ConnectError:
        return {"status": "error", "message": "连接失败，无法访问服务器"}
    except httpx.TimeoutException:
        return {"status": "error", "message": "连接超时"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/api/v1/llm/test")
async def test_llm_connection(provider: str, api_key: str, base_url: str, model: str):
    """测试 LLM 连接"""
    import httpx
    try:
        url = base_url.rstrip('/')
        if '/api/paas/v4' in url:
            endpoint = f"{url}/chat/completions"
        elif url.endswith('/v1') or url.endswith('/v4') or '/v1/' in url or '/v4/' in url:
            endpoint = f"{url}/chat/completions"
        else:
            endpoint = f"{url}/v1/chat/completions"
        
        async with httpx.AsyncClient(timeout=15.0) as client:
            response = await client.post(
                endpoint,
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": model, "messages": [{"role": "user", "content": "hi"}], "max_tokens": 5}
            )
            if response.status_code == 200:
                data = response.json()
                return {"status": "ok", "message": "连接成功", "model": data.get("model", model)}
            else:
                return {"status": "error", "message": f"HTTP {response.status_code}: {response.text[:200]}"}
    except httpx.ConnectError:
        return {"status": "error", "message": "连接失败，无法访问服务器"}
    except httpx.TimeoutException:
        return {"status": "error", "message": "连接超时"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/api/v1/llm/test-image")
async def test_image_model(api_key: str, base_url: str, model: str):
    """测试图像生成模型连接"""
    import httpx
    try:
        url = base_url.rstrip('/')
        if url.endswith('/v1') or url.endswith('/v4') or '/v1/' in url or '/v4/' in url:
            endpoint = f"{url}/images/generations"
        else:
            endpoint = f"{url}/v1/images/generations"
        
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.post(
                endpoint,
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": model, "prompt": "a simple test", "n": 1, "size": "256x256"}
            )
            if response.status_code == 200:
                return {"status": "ok", "message": "图像模型连接成功"}
            else:
                return {"status": "error", "message": f"HTTP {response.status_code}: {response.text[:200]}"}
    except httpx.ConnectError:
        return {"status": "error", "message": "连接失败，无法访问服务器"}
    except httpx.TimeoutException:
        return {"status": "error", "message": "连接超时"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/api/v1/llm/test-video")
async def test_video_model(api_key: str, base_url: str, model: str):
    """测试视频生成模型连接"""
    import httpx
    try:
        url = base_url.rstrip('/')
        if url.endswith('/v1') or url.endswith('/v4') or '/v1/' in url or '/v4/' in url:
            endpoint = f"{url}/video/generations"
        else:
            endpoint = f"{url}/v1/video/generations"
        
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.post(
                endpoint,
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": model, "prompt": "a simple test"}
            )
            if response.status_code in [200, 202]:
                return {"status": "ok", "message": "视频模型连接成功"}
            else:
                return {"status": "error", "message": f"HTTP {response.status_code}: {response.text[:200]}"}
    except httpx.ConnectError:
        return {"status": "error", "message": "连接失败，无法访问服务器"}
    except httpx.TimeoutException:
        return {"status": "error", "message": "连接超时"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/api/v1/llm/providers")
async def list_providers():
    router = get_llm_router()
    providers = [{"name": name, "default_model": getattr(provider, 'default_model', 'unknown'), "is_default": name == router.default_provider} for name, provider in router.providers.items()]
    
    db_settings = get_settings_from_db()
    debug_configured = {}
    for name in ["deepseek", "mimo", "agnes", "qwen", "glm", "openai", "custom"]:
        key = db_settings.get(f"{name}_key", "")
        debug_configured[name] = bool(key and not key.startswith("your-"))
    
    return {"providers": providers, "default_provider": router.default_provider, "debug_configured": debug_configured, "debug_registered": list(router.providers.keys())}

@app.post("/api/v1/chat/provider", response_model=ChatResponse)
async def chat_with_provider(request: ChatWithProviderRequest):
    try:
        router = get_llm_router()
        messages = [{"role": "user", "content": request.message}]
        if request.provider and request.provider in router.providers:
            llm_response = await router.chat(messages, provider=request.provider, model=request.model)
            return ChatResponse(response=llm_response.content, confidence=0.95, sources=[f"llm_{llm_response.provider}"])
        return ChatResponse(response=f"[Demo] 收到您的消息：{request.message}", confidence=0.9, sources=["demo"])
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/v1/content/generate", response_model=ContentResponse)
async def generate_content(request: ContentRequest):
    try:
        router = get_llm_router()
        prompt = f"请用{request.style}风格，写一篇关于「{request.topic}」的文章。"
        messages = [{"role": "user", "content": prompt}]
        if len(router.providers) > 0:
            llm_response = await router.chat_with_fallback(messages, max_tokens=2048)
            content = llm_response.content
        else:
            content = f"[Demo] 关于{request.topic}的{request.style}文章..."
        word_count = len(content)
        reading_time = f"{max(1, word_count // 200)}分钟"
        return ContentResponse(content=content, word_count=word_count, reading_time=reading_time)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/v1/content/generate-doc", response_model=DocumentResponse)
async def generate_document(request: DocumentRequest):
    try:
        router = get_llm_router()
        prompt = f"""请为以下主题生成一份专业的演示文稿内容：{request.topic}
要求：
- 风格：{request.style}
- 共{request.pages}页幻灯片
- 请严格用JSON格式返回，格式如下：
{{
  "title": "演示文稿标题",
  "slides": [
    {{"heading": "页面标题", "content": "详细内容，用换行符分隔多个要点，每个要点一行"}},
    ...
  ]
}}
注意：content中每个要点占一行，不要用列表符号，直接写要点内容。内容要专业、有深度。"""
        if request.extra:
            prompt += f"\n补充说明：{request.extra}"
        
        messages = [{"role": "user", "content": prompt}]
        
        if len(router.providers) > 0:
            llm_response = await router.chat_with_fallback(messages, max_tokens=4096)
            content = llm_response.content
        else:
            content = json.dumps({
                "title": request.topic,
                "slides": [
                    {"heading": f"{request.topic} - 第{i+1}部分", "content": f"这是关于{request.topic}的第{i+1}部分内容..."}
                    for i in range(min(request.pages, 5))
                ]
            })
        
        try:
            doc_data = json.loads(content)
        except:
            doc_data = {"title": request.topic, "slides": [{"heading": request.topic, "content": content}]}
        
        generated_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_media")
        os.makedirs(generated_dir, exist_ok=True)
        
        if request.format == "pptx":
            filename = f"{request.topic[:20]}_{uuid.uuid4().hex[:8]}.pptx"
            filepath = os.path.join(generated_dir, filename)
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            themes = {
                "商务": {"primary": RGBColor(0x1A, 0x3C, 0x6E), "accent": RGBColor(0x2E, 0x86, 0xC1), "bg": RGBColor(0xF8, 0xF9, 0xFA)},
                "学术": {"primary": RGBColor(0x2C, 0x3E, 0x50), "accent": RGBColor(0x27, 0xAE, 0x60), "bg": RGBColor(0xEC, 0xF0, 0xF1)},
                "简洁": {"primary": RGBColor(0x33, 0x33, 0x33), "accent": RGBColor(0xE7, 0x4C, 0x3C), "bg": RGBColor(0xFF, 0xFF, 0xFF)},
                "创意": {"primary": RGBColor(0x8E, 0x44, 0xAD), "accent": RGBColor(0xF3, 0x9C, 0x12), "bg": RGBColor(0xFD, 0xF2, 0xE9)},
            }
            theme = themes.get(request.style, themes["商务"])

            blank_layout = prs.slide_layouts[6]
            title_slide = prs.slide_layouts[0]

            slide = prs.slides.add_slide(title_slide)
            if slide.shapes.title:
                tf = slide.shapes.title.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = doc_data.get("title", request.topic)
                p.font.size = Pt(40)
                p.font.color.rgb = theme["primary"]
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = theme["bg"]

            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.2), Inches(4.333), Inches(0.08)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = theme["accent"]
            accent_bar.line.fill.background()

            slides_data = doc_data.get("slides", [])
            for idx, slide_info in enumerate(slides_data):
                slide = prs.slides.add_slide(blank_layout)
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                top_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.12)
                )
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = theme["accent"]
                top_bar.line.fill.background()

                page_num = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.4))
                tf = page_num.text_frame
                p = tf.paragraphs[0]
                p.text = f"{idx + 1}"
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
                p.alignment = PP_ALIGN.RIGHT

                heading_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
                tf = heading_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_info.get("heading", "")
                p.font.size = Pt(32)
                p.font.color.rgb = theme["primary"]
                p.font.bold = True

                underline = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.55), Inches(2.5), Inches(0.05)
                )
                underline.fill.solid()
                underline.fill.fore_color.rgb = theme["accent"]
                underline.line.fill.background()

                content = slide_info.get("content", "")
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(5.0))
                tf = content_box.text_frame
                tf.word_wrap = True

                lines = content.replace("\\n", "\n").split("\n")
                first = True
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    p.space_after = Pt(10)
                    p.line_spacing = Pt(28)

            prs.save(filepath)
        elif request.format == "docx":
            filename = f"{request.topic[:20]}_{uuid.uuid4().hex[:8]}.docx"
            filepath = os.path.join(generated_dir, filename)
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            doc = Document()

            style = doc.styles['Normal']
            style.font.name = 'Microsoft YaHei'
            style.font.size = Pt(12)
            style.paragraph_format.line_spacing = 1.5

            title_para = doc.add_heading(doc_data.get("title", request.topic), level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

            doc.add_page_break()

            for slide_info in doc_data.get("slides", []):
                heading = doc.add_heading(slide_info.get("heading", ""), level=1)
                content = slide_info.get("content", "")
                lines = content.replace("\\n", "\n").split("\n")
                for line in lines:
                    line = line.strip()
                    if line:
                        doc.add_paragraph(line, style='List Bullet')
                doc.add_paragraph()

            doc.save(filepath)
        else:
            filename = f"{request.topic[:20]}_{uuid.uuid4().hex[:8]}.pdf"
            filepath = os.path.join(generated_dir, filename)
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            from reportlab.pdfbase.cidfonts import UnicodeCIDFont
            pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))
            c = canvas.Canvas(filepath, pagesize=A4)
            width, height = A4
            y = height - 50
            c.setFont('STSong-Light', 18)
            c.drawCentredString(width/2, y, doc_data.get("title", request.topic))
            y -= 40
            for slide_info in doc_data.get("slides", []):
                if y < 100:
                    c.showPage()
                    y = height - 50
                c.setFont('STSong-Light', 14)
                c.drawString(50, y, slide_info.get("heading", ""))
                y -= 25
                c.setFont('STSong-Light', 11)
                content = slide_info.get("content", "")
                lines = content.split('\n')
                for line in lines:
                    if y < 50:
                        c.showPage()
                        y = height - 50
                    c.drawString(50, y, line[:80])
                    y -= 18
                y -= 15
            c.save()
        
        return DocumentResponse(status="ok", filename=filename)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/v1/recommendations", response_model=RecommendationResponse)
async def get_recommendations(request: RecommendationRequest):
    from recommendation_engine import recommend
    recs = recommend(request.user_id, request.top_k)
    scores = [round(0.95 - i * 0.05, 2) for i in range(len(recs))]
    return RecommendationResponse(recommendations=recs, scores=scores)

@app.post("/api/v1/ocr")
async def ocr_recognition(file: UploadFile = File(...)):
    from ocr_engine import ocr_recognize as do_ocr
    content = await file.read()
    return do_ocr(content)

@app.post("/api/v1/ocr/enhance")
async def ocr_enhance(text: str = Query(...)):
    from ocr_engine import ai_enhance
    return {"enhanced": ai_enhance(text)}

@app.post("/api/v1/sentiment")
async def sentiment_analysis(text: str):
    return {"sentiment": "positive", "confidence": 0.89, "details": {"positive": 0.89, "neutral": 0.08, "negative": 0.03}}

# ============ LangChain API ============

@app.get("/api/v1/langchain/status")
async def langchain_status():
    try:
        import langchain
        from langchain_engine import get_rag_engine, get_llm
        rag = get_rag_engine()
        return {
            "langchain_version": langchain.__version__,
            "features": ["LLM包装器", "对话记忆", "链式调用", "Agent智能体", "RAG检索增强", "文档处理"],
            "rag_stats": rag.get_stats(),
            "providers": list(get_llm_router().providers.keys())
        }
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/v1/langchain/chat")
async def langchain_chat(request: ChatStreamRequest):
    from langchain_engine import get_llm, ConversationChain
    llm = get_llm(provider=request.provider or "deepseek")
    conv = ConversationChain(llm=llm)
    response = conv.run(request.message)
    return {"response": response}

@app.post("/api/v1/langchain/rag/add")
async def rag_add_text(text: str = Query(""), file: UploadFile = File(None), uploaded_files: List[UploadFile] = File(None)):
    from langchain_engine import get_rag_engine
    rag = get_rag_engine()
    try:
        total_chunks = 0
        all_files = []
        if uploaded_files:
            all_files = uploaded_files
        elif file:
            all_files = [file]
        if all_files:
            import tempfile
            for upload_file in all_files:
                content = await upload_file.read()
                suffix = os.path.splitext(upload_file.filename or ".txt")[1]
                fd, tmp_path = tempfile.mkstemp(suffix=suffix)
                try:
                    with os.fdopen(fd, 'wb') as tmp:
                        tmp.write(content)
                    total_chunks += rag.add_file(tmp_path)
                finally:
                    if os.path.exists(tmp_path):
                        os.unlink(tmp_path)
        elif text:
            total_chunks = rag.add_text(text)
        else:
            raise HTTPException(status_code=400, detail="需要text或files参数")
        return {"status": "ok", "chunks_added": total_chunks, "stats": rag.get_stats()}
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/v1/langchain/rag/query")
async def rag_query(query: str = Query(...), provider: str = Query(None), top_k: int = Query(3), use_ai: bool = Query(False)):
    try:
        from langchain_engine import get_rag_engine
        rag = get_rag_engine()
        results = rag.search(query, top_k)
        if not results:
            return {"answer": "", "sources": []}
        sources = [{"content": r["content"][:200], "score": r["score"]} for r in results]
        answer = ""
        if use_ai:
            try:
                llm = get_llm_router()
                answer = rag.rag_query(query, llm=llm, top_k=top_k)
            except Exception as e:
                answer = "AI生成回答失败，相关片段已列出。"
        return {"answer": answer, "sources": sources}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/v1/langchain/agent/run")
async def agent_run(query: str, provider: str = None):
    from langchain_engine import get_llm, SimpleAgent, CalculatorTool, DateTimeTool
    llm = get_llm(provider=provider)
    agent = SimpleAgent(llm=llm, tools=[CalculatorTool(), DateTimeTool()])
    response = agent.run(query)
    return {"response": response}

@app.post("/api/v1/langchain/tools/calculate")
async def tool_calculate(expression: str):
    from langchain_engine import CalculatorTool
    tool = CalculatorTool()
    result = tool.run(expression)
    return {"result": result}

@app.get("/api/v1/langchain/tools/time")
async def tool_time():
    from langchain_engine import DateTimeTool
    tool = DateTimeTool()
    return {"result": tool.run()}

@app.get("/api/v1/langchain/tools/list")
async def tool_list():
    from langchain_engine import ALL_TOOLS
    return {"tools": [{"name": t.name, "description": t.description} for t in ALL_TOOLS]}

@app.post("/api/v1/langchain/tools/execute")
async def tool_execute(tool_name: str, input: str = ""):
    from langchain_engine import ALL_TOOLS
    tool_map = {t.name: t for t in ALL_TOOLS}
    if tool_name not in tool_map:
        raise HTTPException(status_code=400, detail=f"Unknown tool: {tool_name}")
    result = tool_map[tool_name].run(input)
    return {"tool": tool_name, "input": input, "result": result}

@app.post("/api/v1/codegen/run")
async def codegen_run(path: str):
    import subprocess
    import os
    import platform
    main_py = os.path.join(path, "main.py")
    main_js = os.path.join(path, "main.js")
    index_html = os.path.join(path, "index.html")
    
    if os.path.exists(main_py):
        try:
            result = subprocess.run(f'cd "{path}" && python main.py', shell=True, capture_output=True, text=True, timeout=30)
            return {"output": result.stdout or "", "errors": result.stderr or "", "type": "python"}
        except subprocess.TimeoutExpired:
            return {"output": "", "errors": "运行超时（30秒）", "type": "python"}
        except Exception as e:
            return {"output": "", "errors": str(e), "type": "python"}
    elif os.path.exists(index_html):
        try:
            system = platform.system().lower()
            if system == "windows":
                os.startfile(index_html)
            elif system == "darwin":
                subprocess.Popen(['open', index_html])
            else:
                subprocess.Popen(['xdg-open', index_html])
            return {"output": f"已在浏览器中打开: {index_html}", "errors": "", "type": "html", "path": index_html}
        except Exception as e:
            return {"output": f"自动打开失败，请手动打开:\n{index_html}", "errors": str(e), "type": "html", "path": index_html}
    elif os.path.exists(main_js):
        try:
            result = subprocess.run(f'cd "{path}" && node main.js', shell=True, capture_output=True, text=True, timeout=30)
            return {"output": result.stdout or "", "errors": result.stderr or "", "type": "node"}
        except Exception as e:
            return {"output": "", "errors": str(e), "type": "node"}
    return {"output": "", "errors": f"未找到可运行文件: {path}", "type": "unknown"}

@app.post("/api/v1/codegen/stream")
async def codegen_stream(request: ChatStreamRequest):
    from langchain_engine import get_llm, CodeGeneratorTool
    
    async def generate():
        llm = get_llm(provider=request.provider, max_tokens=4096)
        
        parts = request.message.split("|")
        description = parts[0].strip() if parts else request.message
        lang = parts[1].strip() if len(parts) > 1 else "python"
        
        lang_map = {"python": "Python", "py": "Python", "html": "HTML/CSS/JavaScript", "web": "HTML/CSS/JavaScript", "js": "JavaScript", "javascript": "JavaScript", "node": "Node.js"}
        lang_name = lang_map.get(lang.lower(), lang)
        
        yield f"data: {json.dumps({'type': 'status', 'content': f'正在分析需求: {description}'})}\n\n"
        
        prompt = f"""你是一个高级全栈程序员。请根据以下描述生成一个完整可运行的{lang_name}项目。

描述: {description}

要求:
1. 代码必须完整、可直接运行，不要省略任何部分
2. 包含必要的注释（中文）
3. Python: 只用标准库，不安装额外依赖
4. HTML: 使用内联CSS和JS，一个文件搞定
5. Node.js: 只用内置模块
6. 代码要美观、有良好的格式
7. 只输出代码，不要任何解释文字

代码:"""
        
        from langchain_core.messages import HumanMessage
        yield f"data: {json.dumps({'type': 'status', 'content': f'正在生成{lang_name}代码...'})}\n\n"
        
        try:
            full_code = ""
            # 用同步LLM调用（CodeGeneratorTool内部已有），这里直接调用
            import httpx
            config = {
                "api_key": llm.api_key,
                "base_url": llm.base_url,
                "model": llm.model
            }
            
            async with httpx.AsyncClient(timeout=120.0) as client:
                response = await client.post(
                    llm.chat_endpoint,
                    headers={"Authorization": f"Bearer {config['api_key']}", "Content-Type": "application/json"},
                    json={
                        "model": config["model"],
                        "messages": [{"role": "user", "content": prompt}],
                        "temperature": 0.3,
                        "max_tokens": 4096,
                        "stream": True
                    }
                )
                
                buffer = ""
                async for line in response.aiter_lines():
                    if line.startswith("data: "):
                        data_str = line[6:].strip()
                        if data_str == "[DONE]":
                            break
                        try:
                            data = json.loads(data_str)
                            delta = data["choices"][0].get("delta", {})
                            content = delta.get("content", "")
                            if content:
                                full_code += content
                                yield f"data: {json.dumps({'type': 'code', 'content': content})}\n\n"
                        except Exception:
                            continue
            
            # 清理代码块标记
            full_code = full_code.strip()
            if full_code.startswith("```"):
                lines = full_code.split("\n")
                lines = lines[1:]
                if lines and lines[-1].strip() == "```":
                    lines = lines[:-1]
                full_code = "\n".join(lines)
            
            yield f"data: {json.dumps({'type': 'status', 'content': '代码生成完成，正在保存项目...'})}\n\n"
            
            # 保存文件
            import os
            projects_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_projects")
            project_name = description[:30].replace(" ", "_").replace("/", "_")
            project_dir = os.path.join(projects_dir, project_name)
            os.makedirs(project_dir, exist_ok=True)
            
            if lang.lower() in ["python", "py"]:
                file_path = os.path.join(project_dir, "main.py")
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(full_code)
                with open(os.path.join(project_dir, "README.md"), "w", encoding="utf-8") as f:
                    f.write(f"# {project_name}\n\n{description}\n\n## 运行\n\n```bash\npython main.py\n```\n")
            elif lang.lower() in ["html", "web"]:
                file_path = os.path.join(project_dir, "index.html")
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(full_code)
            elif lang.lower() in ["js", "javascript", "node"]:
                file_path = os.path.join(project_dir, "main.js")
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(full_code)
            
            yield f"data: {json.dumps({'type': 'status', 'content': f'项目已保存到: {project_dir}'})}\n\n"
            
            # 尝试运行
            if lang.lower() in ["python", "py"]:
                yield f"data: {json.dumps({'type': 'status', 'content': '正在尝试运行项目...'})}\n\n"
                import subprocess
                try:
                    # 检查是否包含需要终端的库
                    needs_terminal = any(k in full_code.lower() for k in ["pygame", "tkinter", "input("])
                    is_windows_curses = platform.system().lower() == "windows" and "curses" in full_code.lower()
                    if needs_terminal or is_windows_curses:
                        yield f"data: {json.dumps({'type': 'status', 'content': '代码需要终端/图形环境运行，已保存项目，请手动运行'})}\n\n"
                    else:
                        result = subprocess.run(f'cd "{project_dir}" && python main.py', shell=True, capture_output=True, text=True, timeout=30)
                        if result.stdout:
                            yield f"data: {json.dumps({'type': 'output', 'content': result.stdout[:3000]})}\n\n"
                        if result.returncode != 0 and result.stderr:
                            # 过滤掉一些无害的警告
                            stderr = result.stderr
                            if "ResourceWarning" not in stderr and len(stderr) > 10:
                                yield f"data: {json.dumps({'type': 'output', 'content': '⚠️ ' + stderr[:1000]})}\n\n"
                except subprocess.TimeoutExpired:
                    yield f"data: {json.dumps({'type': 'status', 'content': '运行超时（30秒），程序可能在等待输入'})}\n\n"
                except Exception as e:
                    yield f"data: {json.dumps({'type': 'output', 'content': f'运行提示: {str(e)[:500]}'})}\n\n"
            elif lang.lower() in ["html", "web"]:
                yield f"data: {json.dumps({'type': 'status', 'content': '请用浏览器打开 index.html 查看效果'})}\n\n"
            
            yield f"data: {json.dumps({'type': 'done', 'path': project_dir})}\n\n"
        
        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'content': str(e)})}\n\n"
        
        yield "data: [DONE]\n\n"
    
    return StreamingResponse(generate(), media_type="text/event-stream", headers={
        "Cache-Control": "no-cache", "Connection": "keep-alive", "X-Accel-Buffering": "no"
    })

@app.post("/api/v1/langchain/agent/chat")
async def agent_chat(request: ChatStreamRequest):
    from langchain_engine import get_llm, SimpleAgent, ALL_TOOLS
    llm = get_llm(provider=request.provider)
    agent = SimpleAgent(llm=llm, tools=ALL_TOOLS)
    response = agent.run(request.message)
    return {"response": response, "tools_used": [t.name for t in agent.tools]}

# ============ 图像和视频生成 ============

class MediaGenerateRequest(BaseModel):
    prompt: str
    provider: str = None
    model: str = None
    negative_prompt: Optional[str] = None
    width: Optional[int] = None
    height: Optional[int] = None
    size: Optional[str] = None
    num: Optional[int] = None
    quality: Optional[str] = None
    style: Optional[str] = None
    seed: Optional[int] = None
    guidance_scale: Optional[float] = None
    duration: Optional[int] = None
    fps: Optional[int] = None
    resolution: Optional[str] = None

@app.post("/api/v1/media/generate-image")
async def generate_image(request: MediaGenerateRequest):
    """生成图像 - 从数据库读取所有配置"""
    import httpx
    settings_db = get_settings_from_db()
    
    img_provider = settings_db.get("image_provider", "agnes")
    
    if img_provider == "custom" and settings_db.get("custom_image_key"):
        api_key = settings_db["custom_image_key"]
        base_url = settings_db.get("custom_image_url", "")
        model = request.model or settings_db.get("custom_image_model", "")
    else:
        api_key = settings_db.get("agnes_key", "")
        base_url = settings_db.get("agnes_url", "https://apihub.agnes-ai.com")
        model = request.model or settings_db.get("agnes_image_model_1", "agnes-image-2.1-flash")
    
    if not api_key:
        raise HTTPException(status_code=400, detail="未配置图像生成 API Key")
    
    img_size = request.size or settings_db.get("image_size", "1024")
    img_num = request.num or int(settings_db.get("image_num", 1))
    img_quality = request.quality or settings_db.get("image_quality", "fhd")
    img_negative = request.negative_prompt or settings_db.get("image_negative_prompt", "")
    img_seed = request.seed if request.seed is not None else settings_db.get("image_seed")
    img_guidance = request.guidance_scale if request.guidance_scale is not None else settings_db.get("image_guidance_scale")
    
    width, height = img_size, img_size
    if "x" in str(img_size):
        parts = str(img_size).split("x")
        width, height = int(parts[0]), int(parts[1])
    
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            payload = {
                "model": model,
                "prompt": request.prompt,
                "n": img_num,
                "size": f"{width}x{height}",
                "quality": img_quality,
            }
            if img_negative:
                payload["negative_prompt"] = img_negative
            if img_seed is not None:
                try: payload["seed"] = int(img_seed)
                except: pass
            if img_guidance is not None:
                try: payload["guidance_scale"] = float(img_guidance)
                except: pass
            
            url = base_url.rstrip('/')
            if url.endswith('/v1') or url.endswith('/v4') or '/v1/' in url or '/v4/' in url:
                img_endpoint = f"{url}/images/generations"
            else:
                img_endpoint = f"{url}/v1/images/generations"
            
            response = await client.post(
                img_endpoint,
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json=payload
            )
            response.raise_for_status()
            data = response.json()
            
            import os
            import base64
            media_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_media", "images")
            os.makedirs(media_dir, exist_ok=True)
            
            saved_files = []
            for i, img in enumerate(data.get("data", [])):
                b64 = img.get("b64_json")
                url = img.get("url")
                if b64:
                    img_bytes = base64.b64decode(b64)
                    filename = f"img_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{i}.png"
                    filepath = os.path.join(media_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(img_bytes)
                    saved_files.append({"filename": filename, "path": filepath})
                elif url:
                    saved_files.append({"url": url})
            
            if not saved_files:
                return {"status": "error", "message": f"API返回格式异常: {json.dumps(data, ensure_ascii=False)[:500]}"}
            
            return {"status": "ok", "images": saved_files, "model": model, "provider": img_provider}
    except httpx.HTTPStatusError as e:
        return {"status": "error", "message": f"API错误: {e.response.status_code} - {e.response.text[:300]}"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/api/v1/media/edit-image")
async def edit_image(
    file: UploadFile = File(...),
    prompt: str = Form(...),
    model: str = Form(None),
    size: str = Form(None),
    num: int = Form(None),
    negative_prompt: str = Form(None),
    seed: int = Form(None),
    guidance_scale: float = Form(None),
):
    """图生图 - 基于参考图编辑生成"""
    import httpx
    import base64
    settings_db = get_settings_from_db()
    
    api_key = settings_db.get("agnes_key", "")
    base_url = settings_db.get("agnes_url", "https://apihub.agnes-ai.com")
    model_name = "agnes-image-2.1-flash"
    
    if not api_key:
        raise HTTPException(status_code=400, detail="未配置图像生成 API Key")
    
    img_size = size or settings_db.get("image_size", "original")
    img_num = num or int(settings_db.get("image_num", 1))
    img_negative = negative_prompt or settings_db.get("image_negative_prompt", "")
    img_seed = seed if seed is not None else settings_db.get("image_seed")
    img_guidance = guidance_scale if guidance_scale is not None else settings_db.get("image_guidance_scale")
    
    try:
        file_content = await file.read()
        file_base64 = base64.b64encode(file_content).decode()
        image_data_uri = f"data:{file.content_type};base64,{file_base64}"
        
        if img_size == "original":
            from PIL import Image
            import io
            img = Image.open(io.BytesIO(file_content))
            width, height = img.size
        else:
            width, height = img_size, img_size
            if "x" in str(img_size):
                parts = str(img_size).split("x")
                width, height = int(parts[0]), int(parts[1])
        
        async with httpx.AsyncClient(timeout=120.0) as client:
            payload = {
                "model": model_name,
                "prompt": prompt,
                "n": img_num,
                "size": f"{width}x{height}",
                "extra_body": {
                    "image": [image_data_uri],
                    "response_format": "b64_json",
                },
            }
            if img_negative:
                payload["negative_prompt"] = img_negative
            if img_seed is not None:
                try: payload["seed"] = int(img_seed)
                except: pass
            if img_guidance is not None:
                try: payload["guidance_scale"] = float(img_guidance)
                except: pass
            
            url = base_url.rstrip('/')
            if url.endswith('/v1') or url.endswith('/v4') or '/v1/' in url or '/v4/' in url:
                img_endpoint = f"{url}/images/generations"
            else:
                img_endpoint = f"{url}/v1/images/generations"
            
            response = await client.post(
                img_endpoint,
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json=payload
            )
            response.raise_for_status()
            data = response.json()
            
            import os
            media_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_media", "images")
            os.makedirs(media_dir, exist_ok=True)
            
            saved_files = []
            for i, img in enumerate(data.get("data", [])):
                b64 = img.get("b64_json")
                url = img.get("url")
                if b64:
                    img_bytes = base64.b64decode(b64)
                    filename = f"img_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{i}.png"
                    filepath = os.path.join(media_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(img_bytes)
                    saved_files.append({"filename": filename, "path": filepath})
                elif url:
                    saved_files.append({"url": url})
            
            if not saved_files:
                return {"status": "error", "message": f"API返回格式异常: {json.dumps(data, ensure_ascii=False)[:500]}"}
            
            return {"status": "ok", "images": saved_files, "model": model_name, "provider": "agnes"}
    except httpx.HTTPStatusError as e:
        return {"status": "error", "message": f"API错误: {e.response.status_code} - {e.response.text[:300]}"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/api/v1/media/generate-video")
async def generate_video(request: MediaGenerateRequest):
    """提交视频生成任务 - 立即返回 task_id"""
    import httpx
    settings_db = get_settings_from_db()
    
    vid_provider = settings_db.get("video_provider", "agnes")
    
    if vid_provider == "custom" and settings_db.get("custom_video_key"):
        api_key = settings_db["custom_video_key"]
        base_url = settings_db.get("custom_video_url", "")
        model = request.model or settings_db.get("custom_video_model", "")
    else:
        api_key = settings_db.get("agnes_key", "")
        base_url = settings_db.get("agnes_url", "https://apihub.agnes-ai.com")
        model = request.model or settings_db.get("agnes_video_model", "agnes-video-v2.0")
    
    if not api_key:
        raise HTTPException(status_code=400, detail="未配置视频生成 API Key")
    
    vid_duration = request.duration or int(settings_db.get("video_duration", 5))
    vid_fps = request.fps or int(settings_db.get("video_fps", 24))
    vid_resolution = request.resolution or settings_db.get("video_resolution", "720p")
    vid_negative = request.negative_prompt or settings_db.get("video_negative_prompt", "")
    vid_seed = request.seed if request.seed is not None else settings_db.get("video_seed")
    
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            payload = {
                "model": model,
                "prompt": request.prompt,
                "duration": vid_duration,
                "fps": vid_fps,
                "resolution": vid_resolution,
            }
            if vid_negative:
                payload["negative_prompt"] = vid_negative
            if vid_seed is not None:
                try: payload["seed"] = int(vid_seed)
                except: pass
            
            url = base_url.rstrip('/')
            if url.endswith('/v1') or url.endswith('/v4') or '/v1/' in url or '/v4/' in url:
                vid_endpoint = f"{url}/video/generations"
            else:
                vid_endpoint = f"{url}/v1/video/generations"
            
            response = await client.post(
                vid_endpoint,
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json=payload
            )
            response.raise_for_status()
            data = response.json()
            
            task_id = data.get("id") or data.get("task_id")
            if not task_id:
                return {"status": "error", "message": "未获取到任务ID", "raw": data}
            
            return {"status": "pending", "message": "任务已提交", "task_id": task_id, "model": model, "provider": vid_provider}
    except httpx.HTTPStatusError as e:
        return {"status": "error", "message": f"API错误: {e.response.status_code} - {e.response.text[:300]}"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/api/v1/media/video-status")
async def video_status(task_id: str):
    """查询视频生成任务状态"""
    import httpx
    import os
    settings_db = get_settings_from_db()
    
    vid_provider = settings_db.get("video_provider", "agnes")
    if vid_provider == "custom" and settings_db.get("custom_video_key"):
        api_key = settings_db["custom_video_key"]
        base_url = settings_db.get("custom_video_url", "")
    else:
        api_key = settings_db.get("agnes_key", "")
        base_url = settings_db.get("agnes_url", "https://apihub.agnes-ai.com")
    
    if not api_key:
        return {"status": "error", "message": "未配置 API Key"}
    
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            url = base_url.rstrip('/')
            if url.endswith('/v1') or url.endswith('/v4') or '/v1/' in url or '/v4/' in url:
                vid_endpoint = f"{url}/video/generations/{task_id}"
            else:
                vid_endpoint = f"{url}/v1/video/generations/{task_id}"
            
            resp = await client.get(
                vid_endpoint,
                headers={"Authorization": f"Bearer {api_key}"}
            )
            resp.raise_for_status()
            status_data = resp.json()
            status = status_data.get("status", "")
            
            if status in ["succeeded", "completed", "done"]:
                media_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_media", "videos")
                os.makedirs(media_dir, exist_ok=True)
                
                video_data = status_data.get("output", status_data.get("data", {}))
                video_url = video_data.get("video_url", "") if isinstance(video_data, dict) else ""
                
                if video_url:
                    video_resp = await client.get(video_url)
                    filename = f"vid_{datetime.now().strftime('%Y%m%d_%H%M%S')}.mp4"
                    filepath = os.path.join(media_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(video_resp.content)
                    return {"status": "ok", "video": {"filename": filename, "path": filepath}}
                
                return {"status": "ok", "video": video_data}
            elif status in ["failed", "error"]:
                return {"status": "error", "message": status_data.get("error", "生成失败")}
            else:
                return {"status": "pending", "message": "生成中", "task_id": task_id}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/api/v1/media/gallery")
async def media_gallery():
    """获取已生成的媒体列表"""
    import os
    media_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_media")
    images = []
    videos = []
    
    img_dir = os.path.join(media_dir, "images")
    if os.path.exists(img_dir):
        for f in sorted(os.listdir(img_dir), reverse=True)[:20]:
            images.append({"filename": f, "path": os.path.join(img_dir, f)})
    
    vid_dir = os.path.join(media_dir, "videos")
    if os.path.exists(vid_dir):
        for f in sorted(os.listdir(vid_dir), reverse=True)[:20]:
            videos.append({"filename": f, "path": os.path.join(vid_dir, f)})
    
    return {"images": images, "videos": videos}

@app.get("/api/v1/media/models")
async def media_models():
    """获取可用的媒体生成模型 - 从数据库读取"""
    settings_db = get_settings_from_db()
    img_provider = settings_db.get("image_provider", "agnes")
    vid_provider = settings_db.get("video_provider", "agnes")
    
    image_models = []
    if img_provider == "custom":
        custom_model = settings_db.get("custom_image_model", "")
        if custom_model:
            image_models.append({"id": custom_model, "name": custom_model, "provider": "custom"})
    else:
        m1 = settings_db.get("agnes_image_model_1", "agnes-image-2.0-flash") or "agnes-image-2.0-flash"
        m2 = settings_db.get("agnes_image_model_2", "agnes-image-2.1-flash") or "agnes-image-2.1-flash"
        image_models.append({"id": m2, "name": m2, "provider": "agnes"})
        if m1 and m1 != m2:
            image_models.append({"id": m1, "name": m1, "provider": "agnes"})
        if not image_models:
            image_models = [{"id": "agnes-image-2.1-flash", "name": "agnes-image-2.1-flash", "provider": "agnes"}]
    
    video_models = []
    if vid_provider == "custom":
        custom_model = settings_db.get("custom_video_model", "")
        if custom_model:
            video_models.append({"id": custom_model, "name": custom_model, "provider": "custom"})
    else:
        m = settings_db.get("agnes_video_model", "agnes-video-v2.0")
        if m: video_models.append({"id": m, "name": m, "provider": "agnes"})
    
    if not image_models:
        image_models = [{"id": "agnes-image-2.1-flash", "name": "agnes-image-2.1-flash", "provider": "agnes"}]
    if not video_models:
        video_models = [{"id": "agnes-video-v2.0", "name": "agnes-video-v2.0", "provider": "agnes"}]
    
    return {
        "image_models": image_models,
        "video_models": video_models,
        "image_provider": img_provider,
        "video_provider": vid_provider,
    }

@app.get("/api/v1/media/preview-doc")
async def preview_document(filename: str):
    """预览文档，生成HTML预览"""
    media_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_media")
    filepath = os.path.join(media_dir, filename)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="文件不存在")

# ============ 工作流 API ============

@app.get("/api/v1/workflows")
async def list_workflows():
    from workflow_engine import list_workflows as do_list
    return {"workflows": do_list()}

@app.post("/api/v1/workflows")
async def create_workflow(data: dict):
    from workflow_engine import create_workflow as do_create
    return do_create(data)

@app.get("/api/v1/workflows/{workflow_id}")
async def get_workflow(workflow_id: str):
    from workflow_engine import get_workflow as do_get
    wf = do_get(workflow_id)
    if not wf:
        raise HTTPException(status_code=404, detail="工作流不存在")
    return wf

@app.put("/api/v1/workflows/{workflow_id}")
async def update_workflow(workflow_id: str, data: dict):
    from workflow_engine import update_workflow as do_update
    ok = do_update(workflow_id, data)
    if not ok:
        raise HTTPException(status_code=404, detail="工作流不存在")
    return {"status": "ok"}

@app.delete("/api/v1/workflows/{workflow_id}")
async def delete_workflow(workflow_id: str):
    from workflow_engine import delete_workflow as do_delete
    ok = do_delete(workflow_id)
    if not ok:
        raise HTTPException(status_code=404, detail="工作流不存在")
    return {"status": "ok"}

@app.post("/api/v1/workflows/{workflow_id}/run")
async def run_workflow(workflow_id: str, data: dict = {}):
    from workflow_engine import run_workflow as do_run
    try:
        llm = get_llm_router()
        run = do_run(workflow_id, inputs=data.get("inputs", {}), llm_router=llm)
        return run.model_dump()
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))

@app.post("/api/v1/workflows/{workflow_id}/run-stream")
async def run_workflow_stream(workflow_id: str, data: dict = {}):
    from workflow_engine import run_workflow as do_run, _cancelled_run_ids
    import threading, json as _json, asyncio

    q = asyncio.Queue()
    run_id_holder = [None]

    def on_progress(event):
        if event.get("type") == "start":
            run_id_holder[0] = event.get("node_id")
        try:
            loop.call_soon_threadsafe(q.put_nowait, event)
        except Exception:
            pass

    loop = asyncio.get_event_loop()

    def _run():
        try:
            llm = get_llm_router()
            do_run(workflow_id, inputs=data.get("inputs", {}), llm_router=llm, progress_callback=on_progress)
        except Exception as e:
            try:
                loop.call_soon_threadsafe(q.put_nowait, {"type": "error", "message": str(e)})
            except Exception:
                pass
        finally:
            try:
                loop.call_soon_threadsafe(q.put_nowait, {"type": "_end"})
            except Exception:
                pass

    t = threading.Thread(target=_run, daemon=True)
    t.start()

    async def event_generator():
        try:
            while True:
                try:
                    event = await asyncio.wait_for(q.get(), timeout=1)
                    if event.get("type") == "_end":
                        yield f"data: {_json.dumps({'type': 'end'})}\n\n"
                        break
                    yield f"data: {_json.dumps(event, ensure_ascii=False, default=str)}\n\n"
                except asyncio.TimeoutError:
                    yield f"data: {_json.dumps({'type': 'ping'})}\n\n"
        finally:
            if run_id_holder[0]:
                from workflow_engine import cancel_workflow
                cancel_workflow(run_id_holder[0])

    from fastapi.responses import StreamingResponse
    return StreamingResponse(event_generator(), media_type="text/event-stream", headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})

@app.post("/api/v1/workflow-runs/{run_id}/cancel")
async def cancel_run(run_id: str):
    from workflow_engine import cancel_workflow
    if cancel_workflow(run_id):
        return {"status": "ok", "message": "已发送取消信号"}
    raise HTTPException(status_code=404, detail="未找到运行中的工作流")

@app.get("/api/v1/workflows/{workflow_id}/runs")
async def get_workflow_runs(workflow_id: str):
    from workflow_engine import get_runs
    return {"runs": get_runs(workflow_id)}

@app.get("/api/v1/workflows/{workflow_id}/export")
async def export_workflow(workflow_id: str):
    from workflow_engine import export_workflow as do_export
    data = do_export(workflow_id)
    if not data:
        raise HTTPException(status_code=404, detail="工作流不存在")
    return data

@app.post("/api/v1/workflows/import")
async def import_workflow(data: dict):
    from workflow_engine import import_workflow as do_import
    return do_import(data)

@app.post("/api/v1/workflows/{workflow_id}/versions")
async def save_version(workflow_id: str, data: dict = {}):
    from workflow_engine import save_version as do_save
    result = do_save(workflow_id, comment=data.get("comment", ""))
    if not result:
        raise HTTPException(status_code=404, detail="工作流不存在")
    return result

@app.get("/api/v1/workflows/{workflow_id}/versions")
async def list_versions(workflow_id: str):
    from workflow_engine import list_versions as do_list
    return {"versions": do_list(workflow_id)}

@app.get("/api/v1/workflows/{workflow_id}/versions/{version}")
async def get_version(workflow_id: str, version: int):
    from workflow_engine import get_version as do_get
    v = do_get(workflow_id, version)
    if not v:
        raise HTTPException(status_code=404, detail="版本不存在")
    return v

@app.post("/api/v1/workflows/{workflow_id}/versions/{version}/restore")
async def restore_version(workflow_id: str, version: int):
    from workflow_engine import restore_version as do_restore
    ok = do_restore(workflow_id, version)
    if not ok:
        raise HTTPException(status_code=404, detail="恢复失败")
    return {"status": "ok"}

@app.delete("/api/v1/workflows/{workflow_id}/versions/{version}")
async def delete_version(workflow_id: str, version: int):
    from workflow_engine import delete_version as do_delete
    ok = do_delete(workflow_id, version)
    if not ok:
        raise HTTPException(status_code=404, detail="版本不存在")
    return {"status": "ok"}

# ==================== 触发器 ====================

_triggers_col = None

def _get_triggers_col():
    global _triggers_col
    if _triggers_col is None:
        from pymongo import MongoClient
        MONGO_URL = os.environ.get("MONGODB_URL", "mongodb://ai_mongo:ai_mongo_2024@localhost:27017")
        client = MongoClient(MONGO_URL, serverSelectionTimeoutMS=3000)
        _triggers_col = client["ai_framework"]["workflow_triggers"]
    return _triggers_col


@app.get("/api/v1/workflows/{workflow_id}/triggers")
async def list_triggers(workflow_id: str):
    col = _get_triggers_col()
    triggers = list(col.find({"workflow_id": workflow_id}, {"_id": 0}))
    return {"triggers": triggers}


@app.post("/api/v1/workflows/{workflow_id}/triggers")
async def create_trigger(workflow_id: str, data: dict):
    col = _get_triggers_col()
    import uuid
    trigger = {
        "id": str(uuid.uuid4()),
        "workflow_id": workflow_id,
        "type": data.get("type", "cron"),
        "config": data.get("config", {}),
        "enabled": data.get("enabled", True),
        "created_at": __import__("datetime").datetime.now().isoformat()
    }
    col.insert_one(trigger)
    trigger.pop("_id", None)
    return trigger


@app.delete("/api/v1/workflows/{workflow_id}/triggers/{trigger_id}")
async def delete_trigger(workflow_id: str, trigger_id: str):
    col = _get_triggers_col()
    col.delete_one({"id": trigger_id, "workflow_id": workflow_id})
    return {"status": "ok"}


@app.post("/api/v1/workflows/{workflow_id}/triggers/{trigger_id}/toggle")
async def toggle_trigger(workflow_id: str, trigger_id: str):
    col = _get_triggers_col()
    trigger = col.find_one({"id": trigger_id, "workflow_id": workflow_id})
    if not trigger:
        raise HTTPException(status_code=404, detail="触发器不存在")
    new_enabled = not trigger.get("enabled", True)
    col.update_one({"id": trigger_id}, {"$set": {"enabled": new_enabled}})
    return {"enabled": new_enabled}


@app.post("/api/v1/workflows/{workflow_id}/trigger/webhook")
async def webhook_trigger(workflow_id: str, data: dict = {}):
    from workflow_engine import run_workflow as do_run
    try:
        llm = get_llm_router()
        run = do_run(workflow_id, inputs=data.get("inputs", {}), llm_router=llm)
        return run.model_dump()
    except ValueError as e:
        raise HTTPException(status_code=404, detail=str(e))


# ==================== 执行队列 ====================

_workflow_queue = queue.Queue()
_queue_worker_running = False
_max_concurrent = 3
_running_count = 0
_queue_lock = threading.Lock()


def _queue_worker():
    global _running_count, _queue_worker_running
    _queue_worker_running = True
    while True:
        try:
            item = _workflow_queue.get(timeout=60)
        except queue.Empty:
            continue
        if item is None:
            break
        with _queue_lock:
            while _running_count >= _max_concurrent:
                _queue_lock.release()
                __import__("time").sleep(1)
                _queue_lock.acquire()
            _running_count += 1
        try:
            wf_id = item["workflow_id"]
            inputs = item.get("inputs", {})
            from workflow_engine import run_workflow as do_run
            llm = get_llm_router()
            run = do_run(wf_id, inputs=inputs, llm_router=llm)
            item["result"] = run.model_dump()
            item["status"] = "completed"
        except Exception as e:
            item["result"] = {"error": str(e)}
            item["status"] = "failed"
        finally:
            with _queue_lock:
                _running_count -= 1
            item["event"].set()


# 启动队列工作线程
_queue_thread = threading.Thread(target=_queue_worker, daemon=True)
_queue_thread.start()


@app.post("/api/v1/workflows/{workflow_id}/queue")
async def queue_workflow(workflow_id: str, data: dict = {}):
    import uuid
    item = {
        "id": str(uuid.uuid4()),
        "workflow_id": workflow_id,
        "inputs": data.get("inputs", {}),
        "priority": data.get("priority", "normal"),
        "status": "queued",
        "created_at": __import__("datetime").datetime.now().isoformat(),
        "result": None,
        "event": threading.Event()
    }
    _workflow_queue.put(item)
    item_no_event = {k: v for k, v in item.items() if k != "event"}
    return item_no_event


@app.get("/api/v1/workflow-queue/status")
async def queue_status():
    return {
        "queue_size": _workflow_queue.qsize(),
        "running": _running_count,
        "max_concurrent": _max_concurrent
    }


@app.post("/api/v1/workflow-queue/config")
async def queue_config(data: dict):
    global _max_concurrent
    _max_concurrent = data.get("max_concurrent", 3)
    return {"max_concurrent": _max_concurrent}


# ==================== 工作流变量 API ====================

@app.put("/api/v1/workflows/{workflow_id}/variables")
async def save_variables(workflow_id: str, data: dict):
    from workflow_engine import get_workflow, update_workflow
    wf = get_workflow(workflow_id)
    if not wf:
        raise HTTPException(status_code=404, detail="工作流不存在")
    variables = data.get("variables", [])
    update_workflow(workflow_id, {"variables": variables})
    return {"status": "ok"}


# ==================== AI创建工作流 ====================

@app.post("/api/v1/workflows/ai-generate")
async def ai_generate_workflow(data: dict):
    description = data.get("description", "")
    if not description:
        raise HTTPException(status_code=400, detail="请提供工作流描述")

    try:
        from llm_config import get_llm_router
        llm = get_llm_router()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM配置加载失败: {e}")

    prompt = f"""你是一个工作流设计专家。根据用户描述生成工作流JSON。

可用节点类型:
- start: 输入节点(起点)
- output: 输出节点(终点)
- llm: LLM调用(config: provider, prompt)
- tool: 工具调用(config: tool_name, input_template, method, url, headers, body, timeout, pattern, file_path, encoding, max_size, smtp_host, smtp_port, smtp_user, smtp_pass, to_addr, subject, write_mode, content, ocr_source, top_k)
- condition: 条件分支(config: condition, operator, value; handles: true, false)
- switch: 多分支(config: switch_var, cases, default_case)
- loop: 循环(config: list_var)
- parallel: 并行(config: empty_config)
- delay: 延时(config: seconds)
- retry: 重试(config: max_retries, delay)
- error_handler: 错误处理(config: fallback_output)
- csv_parse: CSV解析(config: delimiter, has_header, input_template)
- data_filter: 数据过滤(config: data_var, field, operator, value)
- data_sort: 数据排序(config: data_var, sort_field, reverse)
- data_merge: 数据合并(config: left_var, right_var, left_key, right_key)
- deduplicate: 去重(config: data_var, field)
- pivot_table: 透视表(config: data_var, group_by, agg_field, agg_func)
- correlation: 相关性(config: x_var, y_var)
- statistics: 统计(config: data, field, operation)
- chart_gen: 图表(config: data, chart_type[bar/line/pie/scatter], title)
- excel_write: Excel写入(config: file_path, data, sheet_name)
- excel_read: Excel读取(config: input_template, sheet_name)
- docx_write: Word写入(config: save_path, title, content)
- docx_read: Word读取(config: file_path)
- pdf_generate: PDF生成(config: input_template, title)
- text_split: 文本分割(config: delimiter, text)
- text_translate: 翻译(config: text, target_lang)
- text_summarize: 摘要(config: input_template, max_length)
- log_analyze: 日志分析(config: file_path, keyword, max_lines)
- backup: 备份(config: source, dest)
- code_exec: 代码执行(config: language, code, timeout)
- json_build: JSON构建(config: fields)
- hash_encode: 哈希(config: algorithm, input_template)
- datetime: 日期时间(config: action[now/timestamp], format)
- encrypt: 加密(config: text, action[encrypt/decrypt], key)
- calendar_event: 日历(config: title, description, start_time, end_time)
- template_render: 模板(config: template, data, save_path)
- sentiment_analysis: 情感分析(config: text, provider)
- approval: 审批(config: timeout)
- ssh_exec: SSH(config: host, port, username, password, command)
- database: 数据库(config: db_type, connection_string, db_name, collection, query)
- file_operation: 文件操作(config: operation[read/write/copy/move/delete/list], path, content, dest, encoding)
- image_gen: 图像生成(config: prompt, size, model)
- markdown_html: Markdown转HTML(config: input_template)
- regex_replace: 正则替换(config: pattern, replacement, input_template)
- math_calc: 数学计算(config: expression)
- uuid_generate: UUID(config: version, count)
- notify: 通知(config: notify_type, webhook_url, content, title)

规则:
1. 必须有且只有一个start节点和一个output节点
2. 每个节点必须有id(简短如n1,n2), type, label(中文), position{{x,y}}, config
3. 边必须有id, source, target
4. 用 {{{{节点id.字段}}}} 引用上游输出，如 {{{{n1.result}}}}, {{{{s1.input}}}}
5. 只输出JSON，不要其他文字

用户描述: {description}

输出JSON:"""

    try:
        r = ""
        async for chunk in llm.chat_stream([{"role": "user", "content": prompt}], provider="agnes", max_tokens=8192, enable_thinking=False):
            if chunk == "[DONE]" or (isinstance(chunk, str) and chunk.startswith('{"type": "reasoning"')):
                continue
            r += chunk
        raw = r.strip()
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"AI生成失败: {e}")

    if not raw:
        raise HTTPException(status_code=500, detail="AI生成失败: LLM返回空内容，请检查API密钥配置")

    try:
        if "```json" in raw:
            raw = raw.split("```json")[1].split("```")[0].strip()
        elif "```" in raw:
            raw = raw.split("```")[1].split("```")[0].strip()
        wf_data = json.loads(raw)
    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail=f"AI返回格式错误: {raw[:200]}")

    from workflow_engine import create_workflow
    result = create_workflow({
        "name": wf_data.get("name", "AI生成的工作流"),
        "description": description,
        "nodes": wf_data.get("nodes", []),
        "edges": wf_data.get("edges", [])
    })
    return result


# ==================== 节点配置推荐 ====================

@app.get("/api/v1/workflows/{workflow_id}/suggest-config/{node_id}")
async def suggest_config(workflow_id: str, node_id: str):
    from workflow_engine import get_workflow
    wf = get_workflow(workflow_id)
    if not wf:
        raise HTTPException(status_code=404, detail="工作流不存在")

    nodes = wf.get("nodes", [])
    edges = wf.get("edges", [])
    node = next((n for n in nodes if n["id"] == node_id), None)
    if not node:
        raise HTTPException(status_code=404, detail="节点不存在")

    # Find upstream nodes
    upstream_ids = [e["source"] for e in edges if e["target"] == node_id]
    upstream_nodes = [n for n in nodes if n["id"] in upstream_ids]

    suggestions = {}
    ntype = node.get("type", "")

    if ntype == "llm":
        uid = upstream_nodes[0]['id'] if upstream_nodes else 'input'
        suggestions["prompt"] = "处理以下输入:\n{{" + uid + ".result}}"
        suggestions["provider"] = "agnes"

    elif ntype == "chart_gen":
        suggestions["chart_type"] = "bar"
        uid = upstream_nodes[0]['id'] if upstream_nodes else 'input'
        suggestions["data"] = "{{" + uid + ".result}}"
        suggestions["title"] = node.get("label", "Chart")

    elif ntype == "statistics":
        suggestions["data"] = ("{{" + upstream_nodes[0]["id"] + ".result}}" if upstream_nodes else "")
        suggestions["operation"] = "mean"

    elif ntype == "text_translate":
        suggestions["text"] = ("{{" + upstream_nodes[0]["id"] + ".result}}" if upstream_nodes else "")
        suggestions["target_lang"] = "English"

    elif ntype == "text_summarize":
        suggestions["input_template"] = ("{{" + upstream_nodes[0]["id"] + ".result}}" if upstream_nodes else "")
        suggestions["max_length"] = 200

    elif ntype == "pdf_generate":
        suggestions["input_template"] = ("{{" + upstream_nodes[0]["id"] + ".result}}" if upstream_nodes else "")
        suggestions["title"] = node.get("label", "Document")

    elif ntype == "docx_write":
        suggestions["content"] = ("{{" + upstream_nodes[0]["id"] + ".result}}" if upstream_nodes else "")
        suggestions["title"] = node.get("label", "Document")

    elif ntype == "json_build":
        if upstream_nodes:
            suggestions["fields"] = json.dumps({upstream_nodes[0]["id"]: "{{" + upstream_nodes[0]["id"] + ".result}}"})

    elif ntype == "condition":
        suggestions["condition"] = (upstream_nodes[0]["id"] + ".result" if upstream_nodes else "")

    elif ntype == "tool":
        suggestions["tool_name"] = "http"
        suggestions["method"] = "GET"
        suggestions["input_template"] = ("{{" + upstream_nodes[0]["id"] + ".result}}" if upstream_nodes else "")

    elif ntype == "hash_encode":
        suggestions["algorithm"] = "md5"
        suggestions["input_template"] = ("{{" + upstream_nodes[0]["id"] + ".result}}" if upstream_nodes else "")

    elif ntype == "code_exec":
        suggestions["language"] = "python"
        suggestions["code"] = "print('hello')"

    return {"suggestions": suggestions, "upstream": [{"id": n["id"], "label": n.get("label", ""), "type": n["type"]} for n in upstream_nodes]}



@app.get("/api/v1/workflow-templates")
async def list_templates():
    templates = [
        {"category": "运维", "name": "系统巡检日报", "desc": "执行命令+日志分析+LLM报告+PDF", "nodes": 9},
        {"category": "运维", "name": "服务健康监控", "desc": "Ping+HTTP+日志+条件+告警", "nodes": 8},
        {"category": "运维", "name": "日志智能分析", "desc": "扫描+正则+统计+LLM+Word+PDF", "nodes": 8},
        {"category": "数据", "name": "CSV数据分析", "desc": "解析+过滤+排序+统计+图表", "nodes": 6},
        {"category": "办公", "name": "周报生成器", "desc": "分割+LLM润色+PDF+HTML", "nodes": 7},
        {"category": "AI", "name": "RAG知识问答", "desc": "知识库+LLM+摘要+UUID+JSON", "nodes": 8},
    ]
    return {"templates": templates}


@app.get("/api/v1/workflow-runs/{run_id}")
async def get_workflow_run(run_id: str):
    from workflow_engine import get_run
    run = get_run(run_id)
    if not run:
        raise HTTPException(status_code=404, detail="执行记录不存在")
    return run
    
    ext = filename.rsplit('.', 1)[-1].lower()
    if ext == 'pdf':
        try:
            import subprocess
            txt_path = filepath.rsplit('.', 1)[0] + '.txt'
            result = subprocess.run(
                ['python', '-c', f'''
import fitz
doc = fitz.open(r"{filepath}")
text = ""
for page in doc:
    text += page.get_text()
with open(r"{txt_path}", "w", encoding="utf-8") as f:
    f.write(text)
'''],
                capture_output=True, timeout=10
            )
            if os.path.exists(txt_path):
                with open(txt_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                os.remove(txt_path)
            else:
                with open(filepath, 'rb') as f:
                    pdf_bytes = f.read()
                text = f"PDF文件大小：{len(pdf_bytes)} 字节\n无法提取文本内容，请下载后查看。"
            
            paragraphs = ''.join(f'<p>{line}</p>' if line.strip() else '' for line in text.split('\n'))
            html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
body {{ font-family: "Microsoft YaHei", sans-serif; margin: 20px; max-width: 800px; margin: 20px auto; background: #fff; padding: 40px; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); }}
p {{ font-size: 14px; line-height: 1.8; color: #555; margin: 4px 0; }}
</style></head><body>{paragraphs}</body></html>'''
            return {"status": "ok", "html": html}
        except Exception as e:
            return {"status": "error", "message": f"预览失败: {str(e)}"}
    
    try:
        if ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            slides_html = []
            for i, slide in enumerate(prs.slides):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                slides_html.append(f'''
                <div class="slide">
                    <div class="slide-number">幻灯片 {i+1}</div>
                    <div class="slide-content">{"<br>".join(slide_content)}</div>
                </div>''')
            html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
body {{ font-family: "Microsoft YaHei", sans-serif; margin: 20px; background: #f5f5f5; }}
.slide {{ background: #fff; margin: 10px 0; padding: 20px; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); }}
.slide-number {{ color: #999; font-size: 12px; margin-bottom: 10px; }}
.slide-content {{ font-size: 16px; line-height: 1.8; white-space: pre-wrap; }}
</style></head><body>{"".join(slides_html)}</body></html>'''
        elif ext == 'docx':
            from docx import Document
            doc = Document(filepath)
            content_html = []
            for para in doc.paragraphs:
                if para.text.strip():
                    if para.style.name.startswith('Heading'):
                        level = para.style.name.replace('Heading ', '')
                        content_html.append(f'<h{level}>{para.text}</h{level}>')
                    else:
                        content_html.append(f'<p>{para.text}</p>')
            html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8"><style>
body {{ font-family: "Microsoft YaHei", sans-serif; margin: 20px; max-width: 800px; margin: 20px auto; background: #fff; padding: 40px; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); }}
h1,h2,h3 {{ color: #333; }}
p {{ font-size: 14px; line-height: 1.8; color: #555; margin: 8px 0; }}
</style></head><body>{"".join(content_html)}</body></html>'''
        else:
            return {"status": "error", "message": "不支持预览此格式"}
        
        return {"status": "ok", "html": html}
    except Exception as e:
        return {"status": "error", "message": f"预览失败: {str(e)}"}

@app.get("/api/v1/media/download")
async def download_media(filename: str):
    """下载生成的媒体文件"""
    import os
    from fastapi.responses import FileResponse
    media_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "generated_media")
    filepath = os.path.join(media_dir, filename)
    if os.path.exists(filepath):
        ext = filename.rsplit('.', 1)[-1].lower()
        media_types = {
            'html': 'text/html',
            'pdf': 'application/pdf',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        }
        media_type = media_types.get(ext, 'application/octet-stream')
        return FileResponse(filepath, filename=filename, media_type=media_type)
    for subdir in ["images", "videos"]:
        filepath = os.path.join(media_dir, subdir, filename)
        if os.path.exists(filepath):
            return FileResponse(filepath, filename=filename, media_type="application/octet-stream")
    raise HTTPException(status_code=404, detail="文件不存在")

@app.post("/api/v1/tools/test-smtp")
async def test_smtp(data: dict):
    import smtplib
    smtp_host = data.get("smtp_host", "")
    smtp_port = int(data.get("smtp_port", 465))
    smtp_user = data.get("smtp_user", "")
    smtp_pass = data.get("smtp_pass", "")
    use_ssl = data.get("use_ssl", True)
    to_addr = data.get("to_addr", "")
    if not smtp_host or not smtp_user or not smtp_pass:
        raise HTTPException(status_code=400, detail="请填写完整的SMTP服务器、账号和授权码")
    try:
        if use_ssl:
            server = smtplib.SMTP_SSL(smtp_host, smtp_port, timeout=10)
        else:
            server = smtplib.SMTP(smtp_host, smtp_port, timeout=10)
            server.starttls()
        server.login(smtp_user, smtp_pass)
        server.quit()
        result = {"status": "success", "message": f"连接成功: {smtp_host}:{smtp_port}"}
        if to_addr:
            result["message"] += f"，可发送至 {to_addr}"
        return result
    except smtplib.SMTPAuthenticationError:
        raise HTTPException(status_code=400, detail="SMTP 认证失败，请检查账号和授权码")
    except smtplib.SMTPConnectError:
        raise HTTPException(status_code=400, detail=f"无法连接 SMTP 服务器 {smtp_host}:{smtp_port}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"连接失败: {str(e)}")

if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
